from __future__ import annotations

import csv
import io
import shutil
import subprocess
import tempfile
import zipfile
from datetime import date, datetime
from email import policy
from email.parser import BytesParser
from html.parser import HTMLParser
from pathlib import Path

from app.config import settings
from app.schemas.workflow import DocumentSegment
from app.services.workflow.segments import build_text_segments
from app.services.workflow.text_utils import merge_text_sources
from app.services.workflow.uploads.models import ParsedUploadDocument
from app.services.workflow.uploads.ocr import workflow_ocr_service

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".json",
    ".jsonl",
    ".xml",
    ".html",
    ".htm",
    ".yaml",
    ".yml",
    ".ini",
    ".cfg",
    ".log",
    ".sql",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".css",
    ".sh",
    ".eml",
    ".rst",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".heic"}
WORD_EXTENSIONS = {".docx", ".doc"}
EXCEL_EXTENSIONS = {".xlsx", ".xlsm", ".xls"}
PRESENTATION_EXTENSIONS = {".pptx"}
ARCHIVE_EXTENSIONS = {".zip"}

TEXT_CHAR_LIMIT = settings.upload_text_char_limit
TABLE_ROW_LIMIT = settings.upload_table_row_limit
TABLE_COLUMN_LIMIT = settings.upload_table_column_limit
PDF_PAGE_LIMIT = settings.upload_pdf_page_limit
TABLE_SEGMENT_ROW_LIMIT = 8
HEADER_SCAN_ROW_LIMIT = 8
LAYOUT_SCAN_ROW_LIMIT = 300
LAYOUT_SCAN_COLUMN_LIMIT = 20

SHEET_TEMPLATE_HINTS = {
    "template",
    "summary",
    "result",
    "report",
    "form",
    "blank",
    "模板",
    "汇总",
    "填报",
    "填表",
    "回填",
    "结果",
    "空白",
}
TITLE_ROW_HINTS = {
    "汇总",
    "披露",
    "模板",
    "台账",
    "报表",
    "数据表",
    "统计表",
    "清单",
    "说明",
}
HEADER_SIGNAL_ALIASES = {
    "指标",
    "结果",
    "年份",
    "月份",
    "期间",
    "单位",
    "scope",
    "scope1",
    "scope2",
    "scope3",
    "排放量",
    "用量",
    "金额",
    "收入",
    "产量",
    "文档",
    "工作表",
    "数值字段",
}
METRIC_LABEL_ALIASES = {
    "sum": {"sum", "total", "合计", "总计", "汇总", "求和", "总和"},
    "avg": {"avg", "average", "mean", "平均", "平均值"},
    "max": {"max", "maximum", "最大", "最高"},
    "min": {"min", "minimum", "最小", "最低"},
    "count": {"count", "数量", "条数", "行数", "计数"},
    "ratio": {"ratio", "占比", "比例"},
    "intensity": {"intensity", "强度", "密度"},
    "yoy": {"yoy", "同比", "yearoveryear"},
    "mom": {"mom", "环比", "monthovermonth"},
}


class _HtmlTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self.parts.append(text)

    def as_text(self) -> str:
        return " ".join(self.parts)


def infer_document_type(filename: str, content_type: str | None) -> str:
    suffix = Path(filename).suffix.lower()
    mime = (content_type or "").lower()

    if suffix == ".pdf" or mime == "application/pdf":
        return "pdf"
    if suffix in IMAGE_EXTENSIONS or mime.startswith("image/"):
        return "image"
    if suffix in EXCEL_EXTENSIONS or "spreadsheet" in mime or "excel" in mime:
        return "excel"
    if suffix in PRESENTATION_EXTENSIONS or "presentation" in mime or "powerpoint" in mime:
        return "presentation"
    if suffix in WORD_EXTENSIONS or "word" in mime or "document" in mime:
        return "word"
    if suffix == ".eml" or mime == "message/rfc822":
        return "email"
    if suffix == ".msg":
        return "email"
    if suffix in ARCHIVE_EXTENSIONS or mime in {"application/zip", "application/x-zip-compressed"}:
        return "archive"
    if suffix in {".md", ".markdown", ".rst"}:
        return "note"
    if suffix in TEXT_EXTENSIONS or mime.startswith("text/") or mime.endswith("+json"):
        return "text"
    return "other"


def parse_uploaded_bytes(filename: str, content_type: str | None, data: bytes) -> ParsedUploadDocument:
    document_type = infer_document_type(filename, content_type)
    suffix = Path(filename).suffix.lower()

    if not data:
        return ParsedUploadDocument(
            document_type=document_type,
            text=f"[空文件] {filename} 没有可读取内容。",
            parser="empty_file",
            notes=["文件为空，未提取到正文内容。"],
        )

    if document_type == "pdf":
        return _parse_pdf(filename, data)
    if document_type == "image":
        return _parse_image(filename, content_type, data)
    if document_type == "excel":
        return _parse_excel(filename, data)
    if document_type == "presentation":
        return _parse_pptx(filename, data)
    if document_type == "word" and suffix == ".docx":
        return _parse_docx(filename, data)
    if suffix == ".doc":
        return _parse_legacy_doc(filename, data)
    if document_type == "email":
        if suffix == ".msg":
            return _parse_msg(filename, data)
        return _parse_email(filename, data)
    if document_type == "archive":
        return _parse_zip(filename, data)
    if suffix == ".csv":
        return _parse_csv(filename, data)
    if suffix in {".html", ".htm"}:
        return _parse_html(filename, data)
    if document_type in {"text", "note"}:
        return _parse_text(filename, data, parser_name="text_parser")
    decoded = _decode_text(data)
    if decoded is not None:
        text, notes = _finalize_text(decoded)
        return ParsedUploadDocument(
            document_type=document_type,
            text=text,
            parser="fallback_text_parser",
            notes=["按纯文本方式读取该文件。", *notes],
            segments=build_text_segments(text, kind="document", label_prefix="回退片段"),
        )

    return ParsedUploadDocument(
        document_type=document_type,
        text=f"[暂未支持] {filename} 可以上传，但当前环境无法直接读取其正文内容。",
        parser="binary_placeholder",
        notes=["该文件已保留元数据，后续可为此格式补充专用 parser。"],
    )


def _parse_image(filename: str, content_type: str | None, data: bytes) -> ParsedUploadDocument:
    ocr_result = workflow_ocr_service.extract_image_text(
        filename=filename,
        mime_type=content_type,
        data=data,
    )
    if not ocr_result:
        notes = ["图片文件已接收，但当前环境没有拿到 OCR 结果。"]
        if workflow_ocr_service.enabled and not workflow_ocr_service.supports_image_input:
            notes.append(f"当前 OCR provider {settings.model_api_ocr_provider_label} 不支持图片输入。")
        elif not workflow_ocr_service.enabled:
            notes.append("如需自动 OCR，请配置 MODEL_API_OCR_KEY（或共用 MODEL_API_KEY）并启用 MODEL_API_OCR_ENABLED。")
        return ParsedUploadDocument(
            document_type="image",
            text=f"[暂未 OCR] {filename} 已上传，但当前环境还没有接入 OCR 解析。",
            parser="image_placeholder",
            notes=notes,
        )

    text, extra_notes = _finalize_text(ocr_result.text)
    return ParsedUploadDocument(
        document_type="image",
        text=text,
        ocr_text=text,
        parser="image_ocr_parser",
        notes=[*ocr_result.notes, *extra_notes],
        segments=list(ocr_result.segments),
    )


def _parse_text(filename: str, data: bytes, parser_name: str) -> ParsedUploadDocument:
    decoded = _decode_text(data)
    if decoded is None:
        return ParsedUploadDocument(
            document_type="text",
            text=f"[读取失败] {filename} 看起来不像可直接解码的文本文件。",
            parser=f"{parser_name}_failed",
            notes=["文本解码失败，请确认文件编码或格式是否正确。"],
        )

    text, notes = _finalize_text(decoded)
    return ParsedUploadDocument(
        document_type="text",
        text=text,
        parser=parser_name,
        notes=notes,
        segments=build_text_segments(
            text,
            kind="paragraph",
            label_prefix=_default_label_prefix(filename, fallback="文本片段"),
        ),
    )


def _default_label_prefix(filename: str, *, fallback: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in {".md", ".markdown", ".rst"}:
        return "笔记片段"
    if suffix in {".json", ".jsonl", ".xml", ".yaml", ".yml"}:
        return "结构片段"
    return fallback


def _parse_html(filename: str, data: bytes) -> ParsedUploadDocument:
    decoded = _decode_text(data)
    if decoded is None:
        return _parse_text(filename, data, parser_name="html_text_fallback")

    parser = _HtmlTextExtractor()
    parser.feed(decoded)
    extracted = parser.as_text() or decoded
    text, notes = _finalize_text(extracted)
    return ParsedUploadDocument(
        document_type="text",
        text=text,
        parser="html_parser",
        notes=notes,
        segments=build_text_segments(text, kind="section", label_prefix="HTML 片段"),
    )


def _parse_csv(filename: str, data: bytes) -> ParsedUploadDocument:
    decoded = _decode_text(data)
    if decoded is None:
        return _parse_text(filename, data, parser_name="csv_text_fallback")

    reader = csv.reader(io.StringIO(decoded))
    lines: list[str] = []
    notes: list[str] = []
    sheet_rows: list[list[object | None]] = []
    for index, row in enumerate(reader, start=1):
        if TABLE_ROW_LIMIT is not None and index > TABLE_ROW_LIMIT:
            notes.append(f"CSV 只保留前 {TABLE_ROW_LIMIT} 行。")
            break
        clipped_row = list(row[:TABLE_COLUMN_LIMIT]) if TABLE_COLUMN_LIMIT is not None else list(row)
        if TABLE_COLUMN_LIMIT is not None and len(row) > TABLE_COLUMN_LIMIT:
            clipped_row.append("...")
        serialized_row = [_serialize_cell_value(cell.strip()) for cell in clipped_row]
        sheet_rows.append(serialized_row)
        lines.append(" | ".join(_display_cell_value(value) for value in serialized_row))

    sheet_title = Path(filename).stem or "CSV"
    structured_sheet = _build_sheet_structure(sheet_title, sheet_rows)
    text, extra_notes = _finalize_text("\n".join(lines) or decoded)
    return ParsedUploadDocument(
        document_type="excel",
        text=text,
        parser="csv_parser",
        notes=["CSV 已提取结构化表格，可用于后续计算和填表 skill。", *notes, *extra_notes],
        tags=["excel", "spreadsheet", "csv"],
        segments=_build_sheet_segments(sheet_title, sheet_rows),
        structured_data={
            "sheetCount": 1,
            "sheets": [structured_sheet],
        },
    )


def _parse_email(filename: str, data: bytes) -> ParsedUploadDocument:
    message = BytesParser(policy=policy.default).parsebytes(data)
    parts: list[str] = []

    if message.get("subject"):
        parts.append(f"主题：{message['subject']}")
    if message.get("from"):
        parts.append(f"发件人：{message['from']}")
    if message.get("to"):
        parts.append(f"收件人：{message['to']}")

    body = _extract_email_body(message)
    if body:
        parts.append("")
        parts.append(body)

    text, notes = _finalize_text("\n".join(parts) or f"[邮件内容为空] {filename}")
    return ParsedUploadDocument(
        document_type="email",
        text=text,
        parser="email_parser",
        notes=notes,
        segments=_build_email_segments(parts, body),
    )


def _extract_email_body(message) -> str:
    if message.is_multipart():
        for part in message.walk():
            if part.get_content_maintype() == "multipart":
                continue
            if part.get_content_type() == "text/plain":
                return part.get_content()
        for part in message.walk():
            if part.get_content_type() == "text/html":
                html = part.get_content()
                parser = _HtmlTextExtractor()
                parser.feed(html)
                return parser.as_text()
        return ""

    content = message.get_content()
    if message.get_content_type() == "text/html":
        parser = _HtmlTextExtractor()
        parser.feed(content)
        return parser.as_text()
    return content


def _build_email_segments(header_parts: list[str], body: str) -> list[DocumentSegment]:
    segments: list[DocumentSegment] = []
    if header_parts:
        segments.append(
            DocumentSegment(
                segmentId="email_header_1",
                kind="email_header",
                label="邮件头",
                text="\n".join(header_parts),
            )
        )
    if body.strip():
        body_segments = build_text_segments(body, kind="email_body", label_prefix="邮件正文")
        for index, segment in enumerate(body_segments, start=len(segments) + 1):
            segments.append(segment.model_copy(update={"segmentId": f"email_body_{index}"}))
    return segments


def _parse_pdf(filename: str, data: bytes) -> ParsedUploadDocument:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ParsedUploadDocument(
            document_type="pdf",
            text=f"[缺少 PDF 解析依赖] {filename} 已上传，但当前环境未安装 pypdf。",
            parser="pdf_dependency_missing",
            notes=["如需直接读取 PDF 正文，请安装 pypdf。"],
        )

    reader = PdfReader(io.BytesIO(data))
    page_texts: dict[int, str] = {}
    notes: list[str] = []
    total_pages = len(reader.pages)

    for index, page in enumerate(reader.pages, start=1):
        if PDF_PAGE_LIMIT is not None and index > PDF_PAGE_LIMIT:
            notes.append(f"PDF 只提取前 {PDF_PAGE_LIMIT} 页。")
            break
        text = (page.extract_text() or "").strip()
        if text:
            page_texts[index] = text

    processed_pages = min(total_pages, PDF_PAGE_LIMIT) if PDF_PAGE_LIMIT is not None else total_pages
    ocr_result = None
    if workflow_ocr_service.enabled and _should_run_pdf_ocr(page_texts, processed_pages):
        ocr_result = workflow_ocr_service.extract_pdf_text(
            filename=filename,
            data=data,
            max_pages=processed_pages,
        )
        if ocr_result:
            page_texts = _merge_pdf_page_texts(page_texts, ocr_result.page_texts)
            notes.extend(ocr_result.notes)
        elif not workflow_ocr_service.supports_file_input:
            notes.append(
                f"当前 OCR provider {settings.model_api_ocr_provider_label} 不支持 PDF 文件输入，已保留原生 PDF 提取结果。"
            )
        else:
            notes.append(f"{settings.model_api_ocr_provider_label} OCR 未返回可读文字，已保留原生 PDF 提取结果。")
    elif not workflow_ocr_service.enabled and not page_texts:
        notes.append("当前环境未启用 OCR，扫描件 PDF 可能提取不完整。")

    combined_text = (
        "\n\n".join(f"[第 {page_number} 页]\n{page_texts[page_number]}" for page_number in sorted(page_texts))
        if page_texts
        else f"[未提取到正文] {filename} 可能是扫描件或图片型 PDF。"
    )
    if not page_texts:
        notes.append("当前 PDF 没有提取到可读文字。")

    text, extra_notes = _finalize_text(combined_text)
    parser_name = "pdf_parser_with_ocr" if ocr_result else "pdf_parser"
    return ParsedUploadDocument(
        document_type="pdf",
        text=text,
        ocr_text=ocr_result.text if ocr_result else None,
        parser=parser_name,
        notes=[*notes, *extra_notes],
        segments=_build_pdf_segments(page_texts),
    )


def _should_run_pdf_ocr(page_texts: dict[int, str], processed_pages: int) -> bool:
    pdf_mode = settings.model_api_ocr_pdf_mode
    if pdf_mode == "off":
        return False
    if pdf_mode == "hybrid":
        return True
    if not page_texts:
        return True

    pages_with_text = len(page_texts)
    average_chars = sum(len(text) for text in page_texts.values()) / max(pages_with_text, 1)
    return pages_with_text < processed_pages or average_chars < 120


def _merge_pdf_page_texts(
    native_pages: dict[int, str],
    ocr_pages: dict[int, str],
) -> dict[int, str]:
    merged_pages: dict[int, str] = {}
    for page_number in sorted(set(native_pages) | set(ocr_pages)):
        merged_text = merge_text_sources(
            native_pages.get(page_number, ""),
            ocr_pages.get(page_number, ""),
            separator_label="OCR 补全",
        )
        if merged_text:
            merged_pages[page_number] = merged_text
    return merged_pages


def _build_pdf_segments(page_texts: dict[int, str]) -> list[DocumentSegment]:
    segments: list[DocumentSegment] = []
    for page_number in sorted(page_texts):
        page_segments = build_text_segments(
            page_texts[page_number],
            kind="page",
            label_prefix=f"第 {page_number} 页",
            page=page_number,
        )
        for page_segment_index, segment in enumerate(page_segments, start=1):
            segments.append(
                segment.model_copy(update={"segmentId": f"page_{page_number}_{page_segment_index}"})
            )
    return segments


def _parse_docx(filename: str, data: bytes) -> ParsedUploadDocument:
    try:
        from docx import Document
    except ImportError:
        return ParsedUploadDocument(
            document_type="word",
            text=f"[缺少 Word 解析依赖] {filename} 已上传，但当前环境未安装 python-docx。",
            parser="docx_dependency_missing",
            notes=["如需直接读取 .docx，请安装 python-docx。"],
        )

    document = Document(io.BytesIO(data))
    lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    text, notes = _finalize_text("\n".join(lines) or f"[未提取到正文] {filename} 没有检测到段落文本。")
    return ParsedUploadDocument(
        document_type="word",
        text=text,
        parser="docx_parser",
        notes=notes,
        segments=_build_docx_segments(lines),
    )


def _parse_legacy_doc(filename: str, data: bytes) -> ParsedUploadDocument:
    """尽量借助 LibreOffice 转换旧版 .doc。"""
    soffice_binary = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice_binary:
        return ParsedUploadDocument(
            document_type="word",
            text=f"[暂未直接解析] {filename} 是旧版 Word 文档，请先转成 .docx 再上传。",
            parser="word_legacy_placeholder",
            notes=["当前环境只内置 .docx 解析，旧版 .doc 需要 LibreOffice 转换后再读取。"],
        )

    with tempfile.TemporaryDirectory() as temp_dir:
        input_path = Path(temp_dir) / filename
        output_dir = Path(temp_dir) / "converted"
        output_dir.mkdir(parents=True, exist_ok=True)
        input_path.write_bytes(data)
        try:
            subprocess.run(
                [
                    soffice_binary,
                    "--headless",
                    "--convert-to",
                    "docx",
                    "--outdir",
                    str(output_dir),
                    str(input_path),
                ],
                check=True,
                capture_output=True,
            )
        except Exception:
            return ParsedUploadDocument(
                document_type="word",
                text=f"[转换失败] {filename} 未能成功转成 .docx，请检查本地 LibreOffice 环境。",
                parser="word_legacy_conversion_failed",
                notes=["LibreOffice 转换旧版 .doc 失败。"],
            )

        converted_path = output_dir / f"{input_path.stem}.docx"
        if not converted_path.exists():
            return ParsedUploadDocument(
                document_type="word",
                text=f"[转换失败] {filename} 未生成可读取的 .docx 文件。",
                parser="word_legacy_conversion_missing",
                notes=["LibreOffice 未产出转换文件。"],
            )

        parsed = _parse_docx(filename, converted_path.read_bytes())
        parsed.parser = "doc_legacy_via_libreoffice"
        parsed.notes = ["已通过 LibreOffice 转换旧版 .doc 后解析。", *parsed.notes]
        return parsed


def _build_docx_segments(lines: list[str]) -> list[DocumentSegment]:
    segments: list[DocumentSegment] = []
    current_section: str | None = None

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        if len(stripped) <= 28 and not any(punctuation in stripped for punctuation in "。.!?；;：:"):
            current_section = stripped

        built_segments = build_text_segments(
            stripped,
            kind="paragraph",
            label_prefix="段落",
            section=current_section,
        )
        for segment in built_segments:
            index = len(segments) + 1
            segments.append(
                segment.model_copy(
                    update={
                        "segmentId": f"paragraph_{index}",
                        "label": f"段落 {index}",
                    }
                )
            )

    return segments


def _parse_pptx(filename: str, data: bytes) -> ParsedUploadDocument:
    try:
        from pptx import Presentation
    except ImportError:
        return ParsedUploadDocument(
            document_type="presentation",
            text=f"[缺少 PPTX 解析依赖] {filename} 已上传，但当前环境未安装 python-pptx。",
            parser="pptx_dependency_missing",
            notes=["如需直接读取 .pptx，请安装 python-pptx。"],
        )

    presentation = Presentation(io.BytesIO(data))
    slide_lines: list[str] = []
    segments: list[DocumentSegment] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                text_parts.append(text)
        notes_text = ""
        if getattr(slide, "has_notes_slide", False):
            try:
                note_parts = []
                for shape in slide.notes_slide.shapes:
                    text = str(getattr(shape, "text", "") or "").strip()
                    if text:
                        note_parts.append(text)
                notes_text = "\n".join(note_parts).strip()
            except Exception:  # pragma: no cover - defensive against malformed pptx
                notes_text = ""

        combined = "\n".join(text_parts + ([f"备注：{notes_text}"] if notes_text else []))
        if not combined.strip():
            continue

        slide_lines.append(f"[第 {slide_index} 页幻灯片]\n{combined}")
        segments.append(
            DocumentSegment(
                segmentId=f"slide_{slide_index}",
                kind="slide",
                label=f"幻灯片 {slide_index}",
                text=combined.strip(),
                page=slide_index,
            )
        )

    text, notes = _finalize_text("\n\n".join(slide_lines) or f"[未提取到正文] {filename} 没有检测到文本内容。")
    return ParsedUploadDocument(
        document_type="presentation",
        text=text,
        parser="pptx_parser",
        notes=notes,
        segments=segments,
    )


def _parse_msg(filename: str, data: bytes) -> ParsedUploadDocument:
    try:
        import extract_msg
    except ImportError:
        return ParsedUploadDocument(
            document_type="email",
            text=f"[缺少 MSG 解析依赖] {filename} 已上传，但当前环境未安装 extract-msg。",
            parser="msg_dependency_missing",
            notes=["如需直接读取 .msg，请安装 extract-msg。"],
        )

    with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as temp_file:
        temp_file.write(data)
        temp_path = Path(temp_file.name)

    try:
        message = extract_msg.Message(str(temp_path))
        parts = []
        if getattr(message, "subject", ""):
            parts.append(f"主题：{message.subject}")
        if getattr(message, "sender", ""):
            parts.append(f"发件人：{message.sender}")
        if getattr(message, "to", ""):
            parts.append(f"收件人：{message.to}")
        body = str(getattr(message, "body", "") or "").strip()
        attachments = []
        for attachment in getattr(message, "attachments", []) or []:
            attachment_name = str(getattr(attachment, "longFilename", "") or getattr(attachment, "shortFilename", "") or "").strip()
            if attachment_name:
                attachments.append(attachment_name)
        if attachments:
            parts.append(f"附件：{'、'.join(attachments[:8])}")
        if body:
            parts.extend(["", body])
        text, notes = _finalize_text("\n".join(parts) or f"[邮件内容为空] {filename}")
        return ParsedUploadDocument(
            document_type="email",
            text=text,
            parser="msg_parser",
            notes=notes,
            segments=_build_email_segments(parts[:4], body),
        )
    finally:
        temp_path.unlink(missing_ok=True)


def _parse_zip(filename: str, data: bytes, *, depth: int = 0) -> ParsedUploadDocument:
    """递归解析 zip 内支持的文件，并保留 bundle 级摘要。"""
    if depth > 1:
        return ParsedUploadDocument(
            document_type="archive",
            text=f"[停止递归] {filename} 已达到压缩包解析深度上限。",
            parser="zip_depth_limit",
            notes=["压缩包递归解析已达到深度上限。"],
            segments=[
                DocumentSegment(
                    segmentId=f"archive_{depth}_1",
                    kind="archive_summary",
                    label="压缩包概览",
                    text=f"{filename} 已达到压缩包递归解析深度上限。",
                )
            ],
        )

    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        infos = [info for info in archive.infolist() if not info.is_dir()]
        total_entries = len(infos)
        total_size = sum(max(0, info.file_size) for info in infos)

        if total_entries > settings.workflow_zip_entry_limit:
            return ParsedUploadDocument(
                document_type="archive",
                text=f"[条目过多] {filename} 包含 {total_entries} 个文件，超过当前上限。",
                parser="zip_entry_limit",
                notes=[f"压缩包条目过多，当前上限为 {settings.workflow_zip_entry_limit} 个。"],
                segments=[
                    DocumentSegment(
                        segmentId=f"archive_{depth}_1",
                        kind="archive_summary",
                        label="压缩包概览",
                        text=f"{filename} 包含 {total_entries} 个文件，超过当前上限。",
                    )
                ],
            )
        if total_size > settings.workflow_zip_total_size_limit_bytes:
            return ParsedUploadDocument(
                document_type="archive",
                text=f"[体积过大] {filename} 解包体积超过限制，当前未继续展开。",
                parser="zip_size_limit",
                notes=[f"压缩包总解包体积超过 {settings.workflow_zip_total_size_limit_bytes} 字节限制。"],
                segments=[
                    DocumentSegment(
                        segmentId=f"archive_{depth}_1",
                        kind="archive_summary",
                        label="压缩包概览",
                        text=f"{filename} 解包体积超过限制，当前未继续展开。",
                    )
                ],
            )

        child_summaries: list[str] = []
        segments: list[DocumentSegment] = [
            DocumentSegment(
                segmentId=f"archive_{depth}_1",
                kind="archive_summary",
                label="压缩包概览",
                text=f"{filename} 包含 {total_entries} 个文件，总体积约 {total_size} 字节。",
            )
        ]
        notes = ["ZIP 已递归解析支持的子文件。"]
        for child_index, info in enumerate(infos, start=1):
            child_name = Path(info.filename).name or info.filename
            if not child_name:
                continue
            child_data = archive.read(info)
            parsed_child = _parse_zip(child_name, child_data, depth=depth + 1) if child_name.lower().endswith(".zip") else parse_uploaded_bytes(child_name, None, child_data)
            child_summary = f"[{child_name}] {parsed_child.text[:260].strip()}"
            child_summaries.append(child_summary)
            for segment in parsed_child.segments[:8]:
                segments.append(
                    segment.model_copy(
                        update={
                            "segmentId": f"archive_{depth}_{child_index}_{segment.segmentId}",
                            "label": f"{child_name} / {segment.label or '片段'}",
                        }
                    )
                )
            notes.extend(f"{child_name}: {note}" for note in parsed_child.notes[:3])

        text, extra_notes = _finalize_text("\n\n".join(child_summaries) or f"[空压缩包] {filename} 没有可解析文件。")
        return ParsedUploadDocument(
            document_type="archive",
            text=text,
            parser="zip_parser",
            notes=[*notes[:12], *extra_notes],
            segments=segments,
        )


def _parse_excel(filename: str, data: bytes) -> ParsedUploadDocument:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ParsedUploadDocument(
            document_type="excel",
            text=f"[缺少 Excel 解析依赖] {filename} 已上传，但当前环境未安装 openpyxl。",
            parser="excel_dependency_missing",
            notes=["如需直接读取 Excel，请安装 openpyxl。"],
        )

    workbook = load_workbook(io.BytesIO(data), data_only=True)
    lines: list[str] = []
    notes: list[str] = []
    structured_sheets: list[dict[str, object]] = []
    segments: list[DocumentSegment] = []

    for sheet in workbook.worksheets:
        lines.append(f"[工作表] {sheet.title}")
        sheet_rows: list[list[object | None]] = []
        for index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if TABLE_ROW_LIMIT is not None and index > TABLE_ROW_LIMIT:
                notes.append(f"工作表 {sheet.title} 只保留前 {TABLE_ROW_LIMIT} 行。")
                break
            sliced_row = row[:TABLE_COLUMN_LIMIT] if TABLE_COLUMN_LIMIT is not None else row
            values = [_serialize_cell_value(value) for value in sliced_row]
            if TABLE_COLUMN_LIMIT is not None and len(row) > TABLE_COLUMN_LIMIT:
                values.append("...")
            sheet_rows.append(values)
            lines.append(" | ".join(_display_cell_value(value) for value in values).strip())
        lines.append("")
        merged_ranges = [{"range": str(cell_range)} for cell_range in getattr(sheet.merged_cells, "ranges", [])]
        structured_sheets.append(
            _build_sheet_structure(
                sheet.title,
                sheet_rows,
                merged_ranges=merged_ranges,
            )
        )
        segments.extend(
            _build_sheet_segments(
                sheet.title,
                sheet_rows,
                merged_ranges=merged_ranges,
            )
        )

    text, extra_notes = _finalize_text("\n".join(lines).strip() or f"[未提取到表格内容] {filename}")
    workbook.close()
    return ParsedUploadDocument(
        document_type="excel",
        text=text,
        parser="excel_parser",
        notes=["Excel 已提取结构化表格，可用于后续计算和填表 skill。", *notes, *extra_notes],
        tags=["excel", "spreadsheet"],
        segments=segments,
        structured_data={
            "sheetCount": len(structured_sheets),
            "sheets": structured_sheets,
        },
    )


def _build_sheet_segments(
    sheet_title: str,
    rows: list[list[object | None]],
    *,
    merged_ranges: list[dict[str, object]] | None = None,
) -> list[DocumentSegment]:
    structure = _build_sheet_structure(sheet_title, rows, merged_ranges=merged_ranges)
    headers = [str(header) for header in structure.get("headers", [])]
    data_rows = list(structure.get("rows", []))
    segments: list[DocumentSegment] = []
    header_row_index = int(structure.get("headerRowIndex") or 1)
    data_row_start = header_row_index + 1

    if headers:
        summary_text = (
            f"工作表 {sheet_title}。"
            f"表头：{'、'.join(headers[:8])}。"
            f"数据行数：{structure.get('rowCount', 0)}。"
            f"{str(structure.get('layoutSummary') or '').strip()}"
        )
        segments.append(
            DocumentSegment(
                segmentId=f"sheet_{len(segments) + 1}",
                kind="sheet_summary",
                label=f"{sheet_title} 概览",
                text=summary_text,
                sheet=sheet_title,
            )
        )

    for start_index in range(0, len(data_rows), TABLE_SEGMENT_ROW_LIMIT):
        block = data_rows[start_index : start_index + TABLE_SEGMENT_ROW_LIMIT]
        if not block:
            continue
        lines = []
        actual_row_start = data_row_start + start_index
        actual_row_end = actual_row_start + len(block) - 1
        for row_number, record in enumerate(block, start=actual_row_start):
            cells = [
                f"{header}={_display_cell_value(record.get(header))}"
                for header in headers
                if _cell_has_value(record.get(header))
            ]
            if cells:
                lines.append(f"第 {row_number} 行：{'；'.join(cells)}")
        text = "\n".join(lines).strip()
        if not text:
            continue
        segments.append(
            DocumentSegment(
                segmentId=f"sheet_{len(segments) + 1}",
                kind="table_row_block",
                label=f"{sheet_title} 行 {actual_row_start}-{actual_row_end}",
                text=text,
                sheet=sheet_title,
                rowStart=actual_row_start,
                rowEnd=actual_row_end,
            )
        )

    return segments


def _serialize_cell_value(value: object | None) -> object | None:
    if value is None:
        return None
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return value
    if isinstance(value, (datetime, date)):
        return value.isoformat()

    text = str(value).strip()
    return text or None


def _display_cell_value(value: object | None) -> str:
    return "" if value is None else str(value).strip()


def _cell_has_value(value: object | None) -> bool:
    return value not in (None, "")


def _build_sheet_structure(
    sheet_title: str,
    rows: list[list[object | None]],
    *,
    merged_ranges: list[dict[str, object]] | None = None,
) -> dict[str, object]:
    if not any(any(_cell_has_value(value) for value in row) for row in rows):
        return {
            "title": sheet_title,
            "headers": [],
            "rows": [],
            "rowCount": 0,
            "numericColumns": [],
            "headerRowIndex": 1,
            "headerConfidence": "low",
            "sheetRole": "unknown",
            "mergedRanges": list(merged_ranges or []),
            "labelAnchors": [],
            "emptyValueZones": [],
            "layoutSummary": "未识别到稳定表头。",
        }

    header_candidate = _detect_sheet_header(rows)
    header_row_index = header_candidate["rowIndex"]
    header_score = header_candidate["score"]
    header_confidence = _score_to_confidence(header_score)
    header_row = rows[header_row_index - 1] if 0 < header_row_index <= len(rows) else []
    headers = _normalize_headers(header_row)
    data_rows: list[dict[str, object | None]] = []

    for row in rows[header_row_index:]:
        if not any(_cell_has_value(value) for value in row):
            continue
        record = {
            header: row[index] if index < len(row) else None
            for index, header in enumerate(headers)
        }
        if any(_cell_has_value(value) for value in record.values()):
            data_rows.append(record)

    numeric_columns = [
        header
        for header in headers
        if any(_parse_numeric_value(record.get(header)) is not None for record in data_rows)
    ]
    label_anchors = _build_label_anchors(
        rows,
        header_row_index=header_row_index,
        header_confidence=header_confidence,
    )
    empty_value_zones = _build_empty_value_zones(rows, label_anchors)
    sheet_role = _infer_sheet_role(
        sheet_title=sheet_title,
        headers=headers,
        numeric_columns=numeric_columns,
        row_count=len(data_rows),
        label_anchors=label_anchors,
    )
    layout_summary = _build_layout_summary(
        header_row_index=header_row_index,
        header_confidence=header_confidence,
        sheet_role=sheet_role,
        label_anchors=label_anchors,
        empty_value_zones=empty_value_zones,
    )

    return {
        "title": sheet_title,
        "headers": headers,
        "rows": data_rows,
        "rowCount": len(data_rows),
        "numericColumns": numeric_columns,
        "headerRowIndex": header_row_index,
        "headerConfidence": header_confidence,
        "sheetRole": sheet_role,
        "mergedRanges": list(merged_ranges or []),
        "labelAnchors": label_anchors,
        "emptyValueZones": empty_value_zones,
        "layoutSummary": layout_summary,
    }


def _detect_sheet_header(rows: list[list[object | None]]) -> dict[str, int]:
    best_candidate: dict[str, int] | None = None
    scan_limit = min(HEADER_SCAN_ROW_LIMIT, len(rows))

    for row_index in range(1, scan_limit + 1):
        row = rows[row_index - 1]
        non_empty_values = [value for value in row if _cell_has_value(value)]
        if not non_empty_values:
            continue

        non_empty_count = len(non_empty_values)
        numeric_count = sum(_parse_numeric_value(value) is not None for value in non_empty_values)
        header_signal_hits = sum(_looks_like_header_cell(value) for value in non_empty_values)
        score = 0
        score += min(non_empty_count, 4)
        score += min(header_signal_hits, 3) * 3
        if non_empty_count == 1:
            score -= 4
        if numeric_count == 0:
            score += 1
        elif numeric_count >= max(1, non_empty_count - 1):
            score -= 3
        if any(_looks_like_title_cell(value) for value in non_empty_values) and non_empty_count <= 2:
            score -= 3

        next_non_empty_row = _find_next_non_empty_row(rows, row_index)
        if next_non_empty_row and sum(_parse_numeric_value(value) is not None for value in next_non_empty_row) >= 1:
            score += 2

        candidate = {
            "rowIndex": row_index,
            "score": score,
            "signalHits": header_signal_hits,
            "nonEmptyCount": non_empty_count,
        }
        if best_candidate is None:
            best_candidate = candidate
            continue

        if (
            candidate["score"],
            candidate["signalHits"],
            candidate["nonEmptyCount"],
            -candidate["rowIndex"],
        ) > (
            best_candidate["score"],
            best_candidate["signalHits"],
            best_candidate["nonEmptyCount"],
            -best_candidate["rowIndex"],
        ):
            best_candidate = candidate

    if best_candidate is not None:
        return {"rowIndex": best_candidate["rowIndex"], "score": best_candidate["score"]}

    for row_index, row in enumerate(rows, start=1):
        if any(_cell_has_value(value) for value in row):
            return {"rowIndex": row_index, "score": 1}
    return {"rowIndex": 1, "score": 0}


def _find_next_non_empty_row(
    rows: list[list[object | None]],
    current_row_index: int,
) -> list[object | None] | None:
    for row in rows[current_row_index:]:
        if any(_cell_has_value(value) for value in row):
            return row
    return None


def _looks_like_header_cell(value: object | None) -> bool:
    token = _normalize_text_token(value)
    if not token:
        return False
    return token in HEADER_SIGNAL_ALIASES or token.startswith("scope")


def _looks_like_title_cell(value: object | None) -> bool:
    text = _display_cell_value(value)
    token = _normalize_text_token(text)
    if not token:
        return False
    if token in SHEET_TEMPLATE_HINTS:
        return True
    if any(hint in text for hint in TITLE_ROW_HINTS):
        return True
    return len(text) >= 8 and any(hint in text for hint in ("汇总", "披露", "模板", "报表"))


def _score_to_confidence(score: int) -> str:
    if score >= 9:
        return "high"
    if score >= 5:
        return "medium"
    return "low"


def _build_label_anchors(
    rows: list[list[object | None]],
    *,
    header_row_index: int,
    header_confidence: str,
) -> list[dict[str, object]]:
    anchors: list[dict[str, object]] = []
    seen_cells: set[str] = set()

    for row_index in range(1, min(len(rows), LAYOUT_SCAN_ROW_LIMIT) + 1):
        row = rows[row_index - 1]
        row_non_empty = sum(_cell_has_value(value) for value in row)
        for column_index in range(1, min(len(row), LAYOUT_SCAN_COLUMN_LIMIT) + 1):
            cell_value = row[column_index - 1]
            metric_name = _resolve_metric_label(cell_value)
            if metric_name is None:
                continue
            if row_index == header_row_index and header_confidence in {"high", "medium"} and row_non_empty >= 2:
                continue

            cell_ref = _cell_coordinate(row_index, column_index)
            if cell_ref in seen_cells:
                continue
            seen_cells.add(cell_ref)
            anchors.append(
                {
                    "cell": cell_ref,
                    "row": row_index,
                    "column": column_index,
                    "label": _display_cell_value(cell_value),
                    "metric": metric_name,
                }
            )
    return anchors


def _resolve_metric_label(value: object | None) -> str | None:
    normalized = _normalize_text_token(value)
    if not normalized:
        return None

    for metric_name, aliases in METRIC_LABEL_ALIASES.items():
        alias_tokens = {_normalize_text_token(metric_name), *(_normalize_text_token(alias) for alias in aliases)}
        if normalized in alias_tokens:
            return metric_name
    return None


def _build_empty_value_zones(
    rows: list[list[object | None]],
    label_anchors: list[dict[str, object]],
) -> list[dict[str, object]]:
    zones: list[dict[str, object]] = []
    seen_cells: set[tuple[str, str]] = set()

    for anchor in label_anchors:
        row_index = int(anchor.get("row") or 0)
        column_index = int(anchor.get("column") or 0)
        if row_index <= 0 or column_index <= 0:
            continue

        candidate_positions = [
            (row_index, column_index + offset, "right", offset)
            for offset in range(1, 4)
        ]
        candidate_positions.extend(
            (row_index + offset, column_index, "down", offset)
            for offset in range(1, 3)
        )

        added_for_anchor = 0
        for target_row, target_col, direction, distance in candidate_positions:
            candidate_value = _cell_value_at(rows, target_row, target_col)
            if _looks_like_header_cell(candidate_value) or _looks_like_title_cell(candidate_value) or _resolve_metric_label(candidate_value):
                continue

            cell_ref = _cell_coordinate(target_row, target_col)
            dedupe_key = (str(anchor.get("cell") or ""), cell_ref)
            if dedupe_key in seen_cells:
                continue
            seen_cells.add(dedupe_key)
            zones.append(
                {
                    "anchorCell": anchor.get("cell"),
                    "metric": anchor.get("metric"),
                    "cell": cell_ref,
                    "row": target_row,
                    "column": target_col,
                    "direction": direction,
                    "distance": distance,
                }
            )
            added_for_anchor += 1
            if added_for_anchor >= 3:
                break
    return zones


def _cell_value_at(
    rows: list[list[object | None]],
    row_index: int,
    column_index: int,
) -> object | None:
    if row_index < 1 or column_index < 1:
        return None
    if row_index > len(rows):
        return None
    row = rows[row_index - 1]
    if column_index > len(row):
        return None
    return row[column_index - 1]


def _infer_sheet_role(
    *,
    sheet_title: str,
    headers: list[str],
    numeric_columns: list[str],
    row_count: int,
    label_anchors: list[dict[str, object]],
) -> str:
    template_score = 0
    source_score = 0
    normalized_headers = {_normalize_text_token(header) for header in headers}
    sheet_title_lower = sheet_title.lower()

    if {"指标", "结果"} <= normalized_headers:
        template_score += 5
    elif "结果" in normalized_headers or "指标" in normalized_headers:
        template_score += 3
    if label_anchors:
        template_score += min(len(label_anchors), 3) * 2
    if row_count <= 3:
        template_score += 1
    if any(hint in sheet_title_lower for hint in SHEET_TEMPLATE_HINTS):
        template_score += 2

    if numeric_columns:
        source_score += 4
        source_score += min(len(numeric_columns), 2)
    if row_count >= 5:
        source_score += 2
    if any(token in normalized_headers for token in {"月份", "年份", "scope", "scope1", "scope2", "scope3", "金额", "产量", "收入"}):
        source_score += 1

    if template_score >= source_score + 2:
        return "template_like"
    if source_score >= template_score + 1:
        return "source_like"
    return "unknown"


def _build_layout_summary(
    *,
    header_row_index: int,
    header_confidence: str,
    sheet_role: str,
    label_anchors: list[dict[str, object]],
    empty_value_zones: list[dict[str, object]],
) -> str:
    return (
        f"识别表头位于第 {header_row_index} 行（{header_confidence}）。"
        f"工作表类型推断为 {sheet_role}。"
        f"检测到 {len(label_anchors)} 个标签锚点和 {len(empty_value_zones)} 个候选值位。"
    )


def _normalize_text_token(value: object | None) -> str:
    return "".join(
        character
        for character in _display_cell_value(value).lower()
        if character.isalnum() or "\u4e00" <= character <= "\u9fff"
    )


def _cell_coordinate(row_index: int, column_index: int) -> str:
    column_name = ""
    current = column_index
    while current > 0:
        current, remainder = divmod(current - 1, 26)
        column_name = chr(65 + remainder) + column_name
    return f"{column_name}{row_index}"


def _normalize_headers(header_row: list[object | None]) -> list[str]:
    headers: list[str] = []

    for index, value in enumerate(header_row, start=1):
        text = _display_cell_value(value)
        normalized = text or f"column_{index}"
        candidate = normalized
        suffix = 2
        while candidate in headers:
            candidate = f"{normalized}_{suffix}"
            suffix += 1
        headers.append(candidate)

    return headers


def _parse_numeric_value(value: object | None) -> float | None:
    if isinstance(value, bool) or value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)

    text = str(value).strip().replace(",", "")
    if not text:
        return None
    if text.endswith("%"):
        try:
            return float(text[:-1]) / 100
        except ValueError:
            return None
    text = text.replace("¥", "").replace("$", "").replace("￥", "")
    try:
        return float(text)
    except ValueError:
        return None


def _decode_text(data: bytes) -> str | None:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "big5", "shift_jis"):
        try:
            decoded = data.decode(encoding)
        except UnicodeDecodeError:
            continue
        if _looks_like_text(decoded):
            return decoded

    try:
        decoded = data.decode("latin-1")
    except UnicodeDecodeError:
        return None

    return decoded if _looks_like_text(decoded) else None


def _looks_like_text(value: str) -> bool:
    if not value:
        return True

    printable = sum(1 for char in value if char.isprintable() or char in "\n\r\t")
    return printable / len(value) >= 0.85


def _finalize_text(text: str) -> tuple[str, list[str]]:
    compact_text = text.strip()
    if not compact_text:
        return "[未提取到正文] 当前文件没有可读文字。", []
    if TEXT_CHAR_LIMIT is None or len(compact_text) <= TEXT_CHAR_LIMIT:
        return compact_text, []

    return compact_text[:TEXT_CHAR_LIMIT] + "\n\n[内容过长，已截断]", [f"正文过长，只保留前 {TEXT_CHAR_LIMIT} 个字符。"]
