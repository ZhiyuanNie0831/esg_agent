import unittest
import base64
import time
import zipfile
from io import BytesIO
from unittest.mock import patch

from fastapi.testclient import TestClient
from docx import Document
from openpyxl import Workbook, load_workbook
from pypdf import PdfWriter

from app.config import settings
from app.main import app
from app.schemas.workflow import DocumentSegment, IntentionAnalysis, WorkflowDocument, WorkflowExecuteRequest, WorkflowRunRequest
from app.services.workflow.errors import WorkflowAgentUnavailableError, WorkflowConfigurationError
from app.services.workflow.input import WorkflowInputService
from app.services.workflow.esg import resolve_esg_report_word_count
from app.services.workflow.segments import serialize_documents_for_task
from app.services.workflow.skills.base import SkillExecutionContext
from app.services.workflow.skills.fill_validator import FillValidatorSkill
from app.services.workflow.table_fill.data_transfer import TableDataTransferExporter
from app.services.workflow.table_fill.workbook_export import TableFillWorkbookExporter
from app.services.workflow import workflow_agent_service
from app.services.workflow.uploads.ocr import OCRExtraction, WorkflowOCRService


class WorkflowServiceTests(unittest.TestCase):
    def setUp(self) -> None:
        self.input_service = WorkflowInputService()
        self.table_fill_review_patcher = patch(
            "app.services.workflow.table_fill.crosscheck.TableFillMappingCrossCheckService.review_mapping_plan",
            return_value={
                "enabled": False,
                "status": "skipped",
                "approved": True,
                "riskLevel": "unknown",
                "blockWrite": False,
                "issues": [],
                "suggestions": [],
                "candidateReviews": [],
            },
        )
        self.analyze_workflow_patcher = patch.object(
            workflow_agent_service._agent_runtime,
            "analyze_workflow",
            return_value=None,
        )
        self.summarize_documents_patcher = patch.object(
            workflow_agent_service._agent_runtime,
            "summarize_documents",
            return_value=None,
        )
        self.revise_documents_patcher = patch.object(
            workflow_agent_service._agent_runtime,
            "revise_documents",
            return_value=None,
        )
        self.write_esg_report_patcher = patch.object(
            workflow_agent_service._agent_runtime,
            "write_esg_report",
            return_value=None,
        )
        self.compose_final_output_patcher = patch.object(
            workflow_agent_service._agent_runtime,
            "compose_final_output",
            return_value=None,
        )
        self.mock_table_fill_review = self.table_fill_review_patcher.start()
        self.mock_analyze_workflow = self.analyze_workflow_patcher.start()
        self.mock_summarize_documents = self.summarize_documents_patcher.start()
        self.mock_revise_documents = self.revise_documents_patcher.start()
        self.mock_write_esg_report = self.write_esg_report_patcher.start()
        self.mock_compose_final_output = self.compose_final_output_patcher.start()
        self.addCleanup(self.table_fill_review_patcher.stop)
        self.addCleanup(self.analyze_workflow_patcher.stop)
        self.addCleanup(self.summarize_documents_patcher.stop)
        self.addCleanup(self.revise_documents_patcher.stop)
        self.addCleanup(self.write_esg_report_patcher.stop)
        self.addCleanup(self.compose_final_output_patcher.stop)

    def test_skill_catalog_contains_builtin_skills(self) -> None:
        catalog = workflow_agent_service.list_skills()

        self.assertGreaterEqual(catalog.total, 10)
        self.assertIn("document_reader", [skill.name for skill in catalog.skills])
        self.assertIn("document_reviser", [skill.name for skill in catalog.skills])
        self.assertIn("esg_material_checker", [skill.name for skill in catalog.skills])
        self.assertIn("esg_standard_selector", [skill.name for skill in catalog.skills])
        self.assertIn("esg_disclosure_matrix_builder", [skill.name for skill in catalog.skills])
        self.assertIn("esg_kpi_extractor", [skill.name for skill in catalog.skills])
        self.assertIn("esg_data_request_builder", [skill.name for skill in catalog.skills])
        self.assertIn("esg_evidence_linker", [skill.name for skill in catalog.skills])
        self.assertIn("esg_report_writer", [skill.name for skill in catalog.skills])
        self.assertIn("spreadsheet_calculator", [skill.name for skill in catalog.skills])
        self.assertIn("excel_role_classifier", [skill.name for skill in catalog.skills])
        self.assertIn("calculation_planner", [skill.name for skill in catalog.skills])
        self.assertIn("table_mapping_preview", [skill.name for skill in catalog.skills])
        self.assertIn("table_filler", [skill.name for skill in catalog.skills])
        self.assertIn("table_data_transfer", [skill.name for skill in catalog.skills])
        self.assertIn("fill_validator", [skill.name for skill in catalog.skills])

    def test_table_data_transfer_crosscheck_can_block_high_risk_plan(self) -> None:
        class BlockingCrossCheck:
            def review_transfer_plan(self, **kwargs):
                return {
                    "enabled": True,
                    "status": "completed",
                    "approved": False,
                    "riskLevel": "high",
                    "blockWrite": True,
                    "issues": ["源表和目标表可能识别反了。"],
                    "suggestions": ["请人工确认后再写入。"],
                }

        exporter = TableDataTransferExporter(crosscheck_service=BlockingCrossCheck())
        source_document = WorkflowDocument(
            name="数据源表格.xlsx",
            type="excel",
            structuredData={
                "sheets": [
                    {
                        "title": "明细",
                        "headers": ["项目", "金额"],
                        "rows": [{"项目": "打车", "金额": 128}],
                        "numericColumns": ["金额"],
                        "rowCount": 1,
                    }
                ]
            },
        )
        target_document = WorkflowDocument(
            name="目标表格.xlsx",
            type="excel",
            structuredData={
                "sheets": [
                    {
                        "title": "目标填报",
                        "headers": ["项目", "金额"],
                        "rows": [],
                        "numericColumns": [],
                        "rowCount": 0,
                    }
                ]
            },
        )

        result = exporter.build_export_bundle(
            task="将数据源表格写入目标表格。",
            source_documents=[source_document],
            target_documents=[target_document],
        )

        self.assertEqual(result["exportFiles"], [])
        self.assertTrue(result["crossCheck"]["blockWrite"])
        self.assertEqual(result["transferStats"]["blockedByCrossCheck"], 1)
        self.assertEqual(result["transferAudit"][0]["status"], "blocked_by_crosscheck")

    def test_table_data_transfer_fills_existing_keyed_blank_rows_before_appending(self) -> None:
        class SkippedCrossCheck:
            def review_transfer_plan(self, **kwargs):
                return {
                    "enabled": False,
                    "status": "skipped",
                    "approved": True,
                    "riskLevel": "unknown",
                    "blockWrite": False,
                    "issues": [],
                    "suggestions": [],
                }

        exporter = TableDataTransferExporter(crosscheck_service=SkippedCrossCheck())
        source_document = WorkflowDocument(
            name="数据源.xlsx",
            type="excel",
            structuredData={
                "sheets": [
                    {
                        "title": "源数据",
                        "headers": ["项目编号", "项目名称", "金额"],
                        "rows": [
                            {"项目编号": "A-001", "项目名称": "差旅费", "金额": 128},
                            {"项目编号": "A-002", "项目名称": "住宿费", "金额": 523.6},
                        ],
                        "numericColumns": ["金额"],
                        "rowCount": 2,
                    }
                ]
            },
        )
        target_document = WorkflowDocument(
            name="目标模板.xlsx",
            type="excel",
            structuredData={
                "sheets": [
                    {
                        "title": "待填模板",
                        "headers": ["项目编号", "项目名称", "金额"],
                        "rows": [
                            {"项目编号": "A-001", "项目名称": "差旅费", "金额": None},
                            {"项目编号": "A-002", "项目名称": "住宿费", "金额": None},
                        ],
                        "numericColumns": [],
                        "rowCount": 2,
                    }
                ]
            },
        )

        result = exporter.build_export_bundle(
            task="请按项目编号把源数据金额填入目标模板已有空白行。",
            source_documents=[source_document],
            target_documents=[target_document],
        )

        self.assertEqual(result["transferStats"]["mode"], "keyed_blank_fill")
        self.assertEqual(result["transferStats"]["written"], 2)
        self.assertEqual(result["transferStats"]["rowsTransferred"], 2)

        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(result["exportFiles"][0]["contentBase64"]))
        )
        target_sheet = workbook["待填模板"]
        self.assertEqual(target_sheet["A2"].value, "A-001")
        self.assertEqual(target_sheet["B2"].value, "差旅费")
        self.assertEqual(target_sheet["C2"].value, 128)
        self.assertEqual(target_sheet["A3"].value, "A-002")
        self.assertEqual(target_sheet["B3"].value, "住宿费")
        self.assertEqual(target_sheet["C3"].value, 523.6)
        self.assertIsNone(target_sheet["A4"].value)
        self.assertEqual(
            {item["cell"] for item in result["transferAudit"] if item["status"] == "written"},
            {"C2", "C3"},
        )

    def test_table_fill_crosscheck_can_block_high_risk_mapping_plan(self) -> None:
        class BlockingFillCrossCheck:
            def review_mapping_plan(self, **kwargs):
                return {
                    "enabled": True,
                    "status": "completed",
                    "approved": False,
                    "riskLevel": "high",
                    "blockWrite": True,
                    "issues": ["候选填位与模板标签语义不一致。"],
                    "suggestions": ["请人工确认目标单元格。"],
                    "candidateReviews": [],
                }

        exporter = TableFillWorkbookExporter(crosscheck_service=BlockingFillCrossCheck())
        template_document = WorkflowDocument(
            name="回填模板.xlsx",
            type="excel",
            structuredData={
                "sheets": [
                    {
                        "title": "汇总表",
                        "headers": ["指标", "结果"],
                        "rows": [{"指标": "sum", "结果": None}],
                        "numericColumns": [],
                        "rowCount": 1,
                    }
                ]
            },
        )

        result = exporter.build_export_bundle(
            task="计算金额合计并填到模板。",
            documents=[template_document],
            headers=["文档", "工作表", "数值字段", "指标", "结果"],
            rows=[
                {
                    "文档": "源数据.xlsx",
                    "工作表": "明细",
                    "数值字段": "金额",
                    "指标": "sum",
                    "结果": 651.6,
                }
            ],
        )

        self.assertEqual(result["exportFiles"], [])
        self.assertTrue(result["crossCheck"]["blockWrite"])
        self.assertEqual(result["fillStats"]["blockedByCrossCheck"], 1)
        self.assertEqual(result["fillAudit"][0]["status"], "blocked_by_crosscheck")

    def test_table_fill_crosscheck_flags_candidate_for_manual_review(self) -> None:
        class FlaggingFillCrossCheck:
            def review_mapping_plan(self, **kwargs):
                return {
                    "enabled": True,
                    "status": "completed",
                    "approved": True,
                    "riskLevel": "medium",
                    "blockWrite": False,
                    "issues": ["部分候选需要人工复核。"],
                    "suggestions": [],
                    "candidateReviews": [
                        {
                            "mappingId": "map_1",
                            "approved": False,
                            "riskLevel": "high",
                            "issue": "模板标签不像金额合计。",
                            "suggestedSheet": "汇总表",
                            "suggestedCell": "B3",
                        }
                    ],
                }

        exporter = TableFillWorkbookExporter(crosscheck_service=FlaggingFillCrossCheck())
        template_document = WorkflowDocument(
            name="回填模板.xlsx",
            type="excel",
            structuredData={
                "sheets": [
                    {
                        "title": "汇总表",
                        "headers": ["指标", "结果"],
                        "rows": [{"指标": "sum", "结果": None}],
                        "numericColumns": [],
                        "rowCount": 1,
                    }
                ]
            },
        )

        preview = exporter.preview_fill_plan(
            task="计算金额合计并填到模板。",
            documents=[template_document],
            rows=[
                {
                    "文档": "源数据.xlsx",
                    "工作表": "明细",
                    "数值字段": "金额",
                    "指标": "sum",
                    "结果": 651.6,
                }
            ],
        )

        candidate = preview["mappingCandidates"][0]
        self.assertEqual(preview["crossCheck"]["status"], "completed")
        self.assertEqual(candidate["status"], "review_flagged")
        self.assertTrue(candidate["requiresConfirmation"])
        self.assertEqual(candidate["reviewSuggestedCell"], "B3")
        self.assertIn("模板标签不像金额合计", candidate["message"])

    def test_plan_reports_missing_documents(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="帮我检查这份报销材料是否齐全",
                documents=[
                    {
                        "name": "invoice_2026_03.pdf",
                        "type": "pdf",
                        "ocrText": "Invoice amount 523.60.",
                        "tags": ["invoice"],
                    }
                ],
                manualConfirm=True,
            )
        )

        self.assertEqual(response.status, "needs_documents")
        self.assertIn("receipt", response.missingDocuments.missingKinds)
        self.assertEqual(response.missingDocuments.readiness, "partial")
        self.assertTrue(response.missingDocuments.advice)

    def test_plan_requires_documents_when_task_references_materials_without_uploads(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请先总结这些材料的重点，并告诉我有没有缺件",
                documents=[],
                manualConfirm=True,
            )
        )

        self.assertEqual(response.status, "needs_documents")
        self.assertTrue(response.intention.documentRequired)
        self.assertIn("general", response.missingDocuments.missingKinds)

    def test_plan_combines_multi_intent_skills_and_manual_confirmation_checkpoint(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="帮我检查这些报销材料，统计数量，并输出一个简短总结",
                documents=[
                    {
                        "name": "invoice_2026_03.pdf",
                        "type": "pdf",
                        "ocrText": "Invoice amount 523.60.",
                        "tags": ["invoice"],
                    },
                    {
                        "name": "receipt_taxi_0325.png",
                        "type": "image",
                        "ocrText": "Taxi receipt. Amount 128.00.",
                        "tags": ["receipt"],
                    },
                ],
                manualConfirm=True,
            )
        )

        self.assertEqual(response.status, "ready_to_execute")
        self.assertIn("count", response.intention.detectedIntentTypes)
        self.assertIn("summarize", response.intention.detectedIntentTypes)
        self.assertIn("document_counter", response.intention.recommendedSkills)
        self.assertIn("document_summarizer", response.intention.recommendedSkills)
        self.assertTrue(any(step.checkpoint == "approval" for step in response.plan))

    def test_plan_inserts_table_mapping_preview_before_confirmation_for_excel_fill(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="读取原始数据.xlsx 的金额，计算合计和平均值，并填到汇总模板.xlsx。",
                documents=[
                    {
                        "name": "原始数据.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 明细\n项目 | 金额\n打车 | 128\n酒店 | 523.6",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "明细",
                                    "headers": ["项目", "金额"],
                                    "rows": [
                                        {"项目": "打车", "金额": 128},
                                        {"项目": "酒店", "金额": 523.6},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["金额"],
                                }
                            ],
                        },
                    },
                    {
                        "name": "汇总模板.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 汇总表\n指标 | 结果\nsum |\navg |",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "汇总表",
                                    "headers": ["指标", "结果"],
                                    "rows": [
                                        {"指标": "sum", "结果": None},
                                        {"指标": "avg", "结果": None},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": [],
                                }
                            ],
                        },
                    },
                ],
                manualConfirm=True,
                agentMode="off",
            )
        )

        self.assertEqual(response.status, "ready_to_execute")
        ordered_steps = [step.skill or step.checkpoint for step in response.plan if step.skill or step.checkpoint]
        self.assertIn("excel_role_classifier", ordered_steps)
        self.assertIn("calculation_planner", ordered_steps)
        self.assertIn("spreadsheet_calculator", ordered_steps)
        self.assertIn("table_mapping_preview", ordered_steps)
        self.assertIn("approval", ordered_steps)
        self.assertIn("table_filler", ordered_steps)
        self.assertIn("fill_validator", ordered_steps)
        self.assertLess(ordered_steps.index("excel_role_classifier"), ordered_steps.index("calculation_planner"))
        self.assertLess(ordered_steps.index("calculation_planner"), ordered_steps.index("spreadsheet_calculator"))
        self.assertLess(ordered_steps.index("spreadsheet_calculator"), ordered_steps.index("table_mapping_preview"))
        self.assertLess(ordered_steps.index("table_mapping_preview"), ordered_steps.index("approval"))
        self.assertLess(ordered_steps.index("approval"), ordered_steps.index("table_filler"))
        self.assertLess(ordered_steps.index("table_filler"), ordered_steps.index("fill_validator"))

    def test_plan_ignores_unregistered_preferred_skill_and_reports_it(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请帮我总结这份说明",
                documents=[
                    {
                        "name": "notes.txt",
                        "type": "text",
                        "contentText": "这是需要整理的说明文稿。",
                    }
                ],
                preferredSkills=["document_summarizer", "not_registered_skill"],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.status, "ready_to_execute")
        self.assertIn("document_summarizer", response.intention.recommendedSkills)
        self.assertNotIn("not_registered_skill", response.intention.recommendedSkills)
        self.assertIn("not_registered_skill", response.intention.unsupportedPreferredSkills)

    def test_plan_uses_agent_analysis_when_available(self) -> None:
        self.mock_analyze_workflow.return_value = {
            "intentType": "summarize",
            "detectedIntentTypes": ["summarize", "review"],
            "confidence": 0.97,
            "documentRequired": True,
            "requiredDocumentKinds": ["report"],
            "recommendedSkills": ["document_reader", "document_summarizer"],
            "notes": ["agent 已根据材料重排技能顺序。"],
        }

        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请阅读这份报告并提炼重点。",
                documents=[
                    {
                        "name": "report.txt",
                        "type": "text",
                        "contentText": "报告包含排放、员工培训和治理委员会的更新。",
                        "tags": ["report"],
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.status, "ready_to_execute")
        self.assertEqual(response.intention.intentType, "summarize")
        self.assertEqual(response.intention.detectedIntentTypes, ["summarize", "review"])
        self.assertEqual(response.intention.recommendedSkills, ["document_reader", "document_summarizer"])
        self.assertTrue(any("agent" in note.lower() for note in response.intention.notes))

    def test_plan_with_agent_mode_off_skips_agent_analysis(self) -> None:
        self.mock_analyze_workflow.return_value = {
            "intentType": "summarize",
            "detectedIntentTypes": ["summarize"],
            "confidence": 0.95,
            "documentRequired": True,
            "requiredDocumentKinds": ["report"],
            "recommendedSkills": ["document_summarizer"],
            "notes": ["this should not be used"],
        }

        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请帮我总结这份说明",
                documents=[
                    {
                        "name": "notes.txt",
                        "type": "text",
                        "contentText": "这是需要整理的说明文稿。",
                    }
                ],
                manualConfirm=False,
                agentMode="off",
            )
        )

        self.mock_analyze_workflow.assert_not_called()
        self.assertEqual(response.status, "ready_to_execute")
        self.assertIn("document_reader", response.intention.recommendedSkills)
        self.assertTrue(any("关闭 agent" in note for note in response.intention.notes))

    def test_plan_rejects_when_agent_and_local_fallback_are_both_disabled(self) -> None:
        with self.assertRaises(WorkflowConfigurationError):
            workflow_agent_service.plan(
                request=WorkflowRunRequest(
                    task="请帮我总结这份说明",
                    documents=[
                        {
                            "name": "notes.txt",
                            "type": "text",
                            "contentText": "这是需要整理的说明文稿。",
                        }
                    ],
                    manualConfirm=False,
                    agentMode="off",
                    localFallbackEnabled=False,
                )
            )

    def test_plan_requires_agent_result_when_local_fallback_disabled(self) -> None:
        with self.assertRaises(WorkflowAgentUnavailableError):
            workflow_agent_service.plan(
                request=WorkflowRunRequest(
                    task="请帮我总结这份说明",
                    documents=[
                        {
                            "name": "notes.txt",
                            "type": "text",
                            "contentText": "这是需要整理的说明文稿。",
                        }
                    ],
                    manualConfirm=False,
                    agentMode="on",
                    localFallbackEnabled=False,
                )
            )

    def test_plan_esg_demo_can_report_missing_contract_attachment(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请帮我检查这套 ESG 披露材料是否齐全，确认是否包含报告正文和供应商合同附件，并输出缺件提示。",
                documents=[
                    {
                        "name": "2025_esg_report_draft.docx",
                        "type": "word",
                        "contentText": "2025 ESG 报告草稿，包含环境、社会与治理章节。",
                        "tags": ["report"],
                    }
                ],
                manualConfirm=True,
            )
        )

        self.assertEqual(response.status, "needs_documents")
        self.assertIn("report", response.missingDocuments.presentKinds)
        self.assertIn("contract", response.missingDocuments.missingKinds)

    def test_plan_esg_task_recommends_specialized_esg_skills(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请检查这套 ESG 披露材料的覆盖度，提取关键 KPI，并给出报告章节大纲。",
                documents=[
                    {
                        "name": "2025_esg_report_draft.docx",
                        "type": "word",
                        "contentText": "环境部分包括可再生电力占比 48%。社会部分包括培训覆盖率 91%。治理部分包括反舞弊培训覆盖率 100%。",
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.status, "ready_to_execute")
        self.assertIn("esg_standard_selector", response.intention.recommendedSkills)
        self.assertIn("esg_disclosure_mapper", response.intention.recommendedSkills)
        self.assertIn("esg_disclosure_matrix_builder", response.intention.recommendedSkills)
        self.assertIn("esg_material_checker", response.intention.recommendedSkills)
        self.assertIn("esg_kpi_extractor", response.intention.recommendedSkills)
        self.assertIn("esg_data_request_builder", response.intention.recommendedSkills)
        self.assertIn("esg_evidence_linker", response.intention.recommendedSkills)
        self.assertIn("esg_report_outline_builder", response.intention.recommendedSkills)
        self.assertIn("要回答的问题", response.plan[0].description)
        self.assertTrue(any("标明来源片段" in step.description for step in response.plan if step.skill == "esg_kpi_extractor"))

    def test_plan_prepared_documents_include_segments_for_manual_text(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请总结这份 ESG 说明",
                documents=[
                    {
                        "name": "esg_note.txt",
                        "type": "text",
                        "contentText": "环境部分提到可再生电力占比 48%。\n\n治理部分提到反舞弊培训覆盖率 100%。",
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.status, "ready_to_execute")
        prepared_document = response.preparedDocuments[0]
        self.assertGreaterEqual(len(prepared_document.segments), 2)
        self.assertIn("可再生电力占比 48%", prepared_document.segments[0].text)

    def test_plan_natural_language_review_request_is_not_classified_as_general(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请帮我看一下这份 ESG 报告，告诉我主要问题。",
                documents=[
                    {
                        "name": "esg_report.txt",
                        "type": "text",
                        "contentText": "报告草稿提到董事会治理、培训覆盖率和碳排放目标。",
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.status, "ready_to_execute")
        self.assertEqual(response.intention.intentType, "review")
        self.assertNotEqual(response.intention.intentType, "general")

    def test_plan_natural_language_revision_request_is_not_classified_as_general(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="把这段治理章节改得更正式一点，整理成可以进报告的版本。",
                documents=[
                    {
                        "name": "governance_draft.txt",
                        "type": "text",
                        "contentText": "治理草稿：董事会监督 ESG，反舞弊培训覆盖率 100%。",
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.intention.intentType, "revise")
        self.assertIn("document_reviser", response.intention.recommendedSkills)

    def test_plan_missing_check_request_recognizes_missing_intent_from_natural_language(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="检查这套 ESG 材料缺什么，是否已经齐全。",
                documents=[
                    {
                        "name": "esg_report.txt",
                        "type": "text",
                        "contentText": "这是 ESG 报告正文草稿。",
                        "tags": ["report"],
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.intention.intentType, "check_missing")
        self.assertIn("check_missing", response.intention.detectedIntentTypes)

    def test_plan_calculation_request_recognizes_count_intent_from_natural_language(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="帮我算一下这个碳排台账里的合计和平均值。",
                documents=[
                    {
                        "name": "emissions.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] Scope 数据",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "Scope 数据",
                                    "headers": ["月份", "Scope 1", "Scope 2"],
                                    "rows": [{"月份": "1月", "Scope 1": 10, "Scope 2": 20}],
                                    "rowCount": 1,
                                    "numericColumns": ["Scope 1", "Scope 2"],
                                }
                            ],
                        },
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.intention.intentType, "count")
        self.assertIn("spreadsheet_calculator", response.intention.recommendedSkills)
        self.assertNotIn("table_filler", response.intention.recommendedSkills)

    def test_plan_result_table_wording_does_not_trigger_template_fill(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请读取这份 Excel，计算金额合计和平均值，并生成一张结果表。",
                documents=[
                    {
                        "name": "费用明细.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 明细\n项目 | 金额\n打车 | 128\n酒店 | 523.6",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "明细",
                                    "headers": ["项目", "金额"],
                                    "rows": [
                                        {"项目": "打车", "金额": 128},
                                        {"项目": "酒店", "金额": 523.6},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["金额"],
                                }
                            ],
                        },
                    }
                ],
                manualConfirm=False,
            )
        )

        skill_names = [step.skill for step in response.plan if step.skill]
        self.assertIn("calculation_planner", skill_names)
        self.assertIn("spreadsheet_calculator", skill_names)
        self.assertNotIn("table_filler", skill_names)
        self.assertFalse(any(step.checkpoint == "approval" for step in response.plan))

    def test_plan_template_fill_requires_mapping_confirmation_without_global_manual_confirm(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="读取原始数据.xlsx 的金额，计算合计和平均值，并填到空白汇总模板.xlsx。",
                documents=[
                    {
                        "name": "原始数据.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 明细\n项目 | 金额\n打车 | 128\n酒店 | 523.6",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "明细",
                                    "headers": ["项目", "金额"],
                                    "rows": [
                                        {"项目": "打车", "金额": 128},
                                        {"项目": "酒店", "金额": 523.6},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["金额"],
                                }
                            ],
                        },
                    },
                    {
                        "name": "空白汇总模板.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 汇总页\n指标 | 结果\nsum | \navg | ",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "汇总页",
                                    "headers": ["指标", "结果"],
                                    "rows": [
                                        {"指标": "sum", "结果": None},
                                        {"指标": "avg", "结果": None},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": [],
                                }
                            ],
                        },
                    },
                ],
                manualConfirm=False,
            )
        )

        ordered_steps = [step.skill or step.checkpoint for step in response.plan]
        self.assertIn("approval", ordered_steps)
        self.assertLess(ordered_steps.index("table_mapping_preview"), ordered_steps.index("approval"))
        self.assertLess(ordered_steps.index("approval"), ordered_steps.index("table_filler"))
        approval_step = next(step for step in response.plan if step.checkpoint == "approval")
        self.assertEqual(approval_step.inputs["confirmationType"], "table_mapping")

    def test_plan_summary_request_recognizes_key_points_question(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="读完之后告诉我重点。",
                documents=[
                    {
                        "name": "note.txt",
                        "type": "text",
                        "contentText": "报告重点包括培训、排放和治理委员会。",
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertEqual(response.intention.intentType, "summarize")
        self.assertNotEqual(response.intention.intentType, "general")

    def test_plan_outline_request_is_not_classified_as_general(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="我想做一个 ESG 报告大纲。",
                documents=[
                    {
                        "name": "esg_notes.txt",
                        "type": "text",
                        "contentText": "材料包括环境、社会和治理三个章节的披露重点。",
                    }
                ],
                manualConfirm=False,
            )
        )

        self.assertNotEqual(response.intention.intentType, "general")
        self.assertIn("esg_report_outline_builder", response.intention.recommendedSkills)
        self.assertNotIn("esg_report_writer", response.intention.recommendedSkills)

    def test_plan_esg_report_request_recommends_report_writer(self) -> None:
        response = workflow_agent_service.plan(
            request=WorkflowRunRequest(
                task="请根据上传材料生成一份 1200 字 ESG 报告。",
                documents=[
                    {
                        "name": "esg_notes.txt",
                        "type": "text",
                        "contentText": "环境部分包括可再生电力占比 48%。社会部分包括培训覆盖率 91%。治理部分包括反舞弊培训覆盖率 100%。",
                    }
                ],
                manualConfirm=False,
                agentMode="off",
            )
        )

        self.assertEqual(response.status, "ready_to_execute")
        self.assertIn("esg_report_outline_builder", response.intention.recommendedSkills)
        self.assertIn("esg_report_writer", response.intention.recommendedSkills)
        ordered_steps = [step.skill or step.checkpoint for step in response.plan if step.skill or step.checkpoint]
        self.assertLess(ordered_steps.index("esg_report_outline_builder"), ordered_steps.index("approval"))
        self.assertLess(ordered_steps.index("approval"), ordered_steps.index("esg_report_writer"))
        writer_step = next(step for step in response.plan if step.skill == "esg_report_writer")
        self.assertIn("客户字数要求", writer_step.description)

    def test_resolve_esg_report_word_count_from_task(self) -> None:
        requirement = resolve_esg_report_word_count("请生成一份 1.5万字 ESG 报告")

        self.assertTrue(requirement["explicit"])
        self.assertEqual(requirement["targetWordCount"], 15000)
        self.assertGreaterEqual(requirement["minWordCount"], 13500)

    def test_segment_selection_prefers_task_relevant_snippets(self) -> None:
        documents = self.input_service.prepare_documents(
            [
                WorkflowDocument(
                    name="climate_note.txt",
                    type="text",
                    contentText="可再生电力占比提升至 48%，范围二排放下降 12%。",
                ),
                WorkflowDocument(
                    name="governance_note.txt",
                    type="text",
                    contentText="Board ESG committee met 4 times. Anti-bribery training coverage reached 100%.",
                ),
            ]
        )

        payload = serialize_documents_for_task(
            "请提炼 anti-bribery 和 governance 风险亮点",
            documents,
            max_documents=2,
            max_segments_per_document=2,
            max_total_segments=4,
            segment_text_limit=240,
        )

        self.assertIn("anti-bribery", payload["taskKeywords"])
        self.assertEqual(len(payload["documents"]), 1)
        self.assertEqual(payload["documents"][0]["name"], "governance_note.txt")
        self.assertIn("Anti-bribery", payload["documents"][0]["selectedSegments"][0]["text"])
        self.assertIn("anti-bribery", payload["documents"][0]["selectedSegments"][0]["matchedKeywords"])

    def test_execute_without_agent_still_uses_task_relevant_segments_for_summary(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="请总结这份材料里的 anti-bribery 风险",
                documents=[
                    {
                        "name": "risk_note.txt",
                        "type": "text",
                        "contentText": (
                            "公司背景说明，介绍业务范围和时间线。"
                            "\n\n环境章节主要介绍能耗、排放和可再生电力采购。"
                            "\n\nAnti-bribery training coverage reached 100%, and no incidents were reported."
                            "\n\n治理部分还提到董事会 ESG 委员会召开 4 次。"
                        ),
                    }
                ],
                manualConfirm=False,
                approved=True,
                agentMode="off",
            )
        )

        self.assertEqual(response.status, "completed")
        self.assertIn("document_summarizer", response.executedSkills)
        self.assertTrue(response.finalOutput)
        self.assertIn("anti-bribery", response.finalOutput.summaryText.lower())

    def test_pdf_ocr_retries_page_by_page_when_multi_page_output_loses_markers(self) -> None:
        buffer = BytesIO()
        writer = PdfWriter()
        writer.add_blank_page(width=72, height=72)
        writer.add_blank_page(width=72, height=72)
        writer.write(buffer)

        service = WorkflowOCRService()
        service._enabled = True
        service._client = type("DummyClient", (), {"supports_file_input": True})()

        original_pages_per_request = settings.model_api_ocr_pdf_pages_per_request
        object.__setattr__(settings, "model_api_ocr_pdf_pages_per_request", 2)
        try:
            with patch.object(
                service,
                "_request_text",
                side_effect=[
                    "这是一段没有页码标记的整块 OCR 输出。",
                    "第一页单独重试结果。",
                    "第二页单独重试结果。",
                ],
            ):
                result = service.extract_pdf_text(
                    filename="scanned.pdf",
                    data=buffer.getvalue(),
                    max_pages=2,
                )
        finally:
            object.__setattr__(settings, "model_api_ocr_pdf_pages_per_request", original_pages_per_request)

        self.assertIsNotNone(result)
        self.assertEqual(result.page_texts[1], "第一页单独重试结果。")
        self.assertEqual(result.page_texts[2], "第二页单独重试结果。")
        self.assertTrue(any("逐页重试" in note for note in result.notes))

    def test_execute_returns_summary_and_revised_document(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="请帮我修订这两份说明，整理成一版更清晰的草稿",
                documents=[
                    {
                        "name": "draft_a.txt",
                        "type": "text",
                        "contentText": "第一版草稿，包含背景、目标和里程碑。",
                    },
                    {
                        "name": "draft_b.txt",
                        "type": "text",
                        "contentText": "第二版补充了风险、资源和交付计划。",
                    },
                ],
                manualConfirm=True,
                approved=True,
            )
        )

        self.assertEqual(response.status, "completed")
        self.assertIn("document_reviser", response.executedSkills)
        self.assertTrue(response.finalOutput)
        self.assertIn("修订稿", response.finalOutput.revisedDocument or "")
        self.assertTrue(response.finalOutput.nextActions)

    def test_execute_esg_management_summary_demo_returns_summary(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="请阅读这三份 ESG 材料，提炼一版给管理层的简短摘要，突出气候、员工与治理亮点。",
                documents=[
                    {
                        "name": "climate_progress_note.txt",
                        "type": "text",
                        "contentText": "可再生电力采购比例提升至 48%，范围二排放同比下降 12%。",
                    },
                    {
                        "name": "people_metrics_update.docx",
                        "type": "word",
                        "contentText": "关键岗位培训覆盖率 91%，女性管理者占比 39%。",
                    },
                    {
                        "name": "governance_committee_brief.pdf",
                        "type": "pdf",
                        "ocrText": "Board ESG committee met 4 times. Anti-bribery training coverage reached 100%.",
                    },
                ],
                manualConfirm=True,
                approved=True,
            )
        )

        self.assertEqual(response.status, "completed")
        self.assertIn("document_summarizer", response.executedSkills)
        self.assertTrue(response.finalOutput)
        self.assertIn("已处理文档数：3", response.finalOutput.summaryText)
        self.assertTrue(response.finalOutput.evidence)
        self.assertTrue(response.finalOutput.evidence[0].document)
        self.assertTrue(response.finalOutput.evidence[0].location)

    def test_execute_esg_skillset_returns_indicators_and_outline(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="请提取这批 ESG 材料里的关键指标，并生成一版披露大纲。",
                documents=[
                    {
                        "name": "esg_kpi_pack.txt",
                        "type": "text",
                        "contentText": (
                            "可再生电力占比提升至 48%。"
                            "关键岗位培训覆盖率 91%。"
                            "Board ESG committee met 4 times."
                            "Anti-bribery training coverage reached 100%."
                        ),
                    }
                ],
                manualConfirm=False,
                approved=True,
            )
        )

        self.assertEqual(response.status, "completed")
        self.assertIn("esg_standard_selector", response.executedSkills)
        self.assertIn("esg_disclosure_matrix_builder", response.executedSkills)
        self.assertIn("esg_kpi_extractor", response.executedSkills)
        self.assertIn("esg_evidence_linker", response.executedSkills)
        self.assertIn("esg_report_outline_builder", response.executedSkills)
        standard_result = response.finalOutput.artifacts["技能结果"]["esg_standard_selector"]
        matrix_result = response.finalOutput.artifacts["技能结果"]["esg_disclosure_matrix_builder"]
        kpi_result = response.finalOutput.artifacts["技能结果"]["esg_kpi_extractor"]
        linker_result = response.finalOutput.artifacts["技能结果"]["esg_evidence_linker"]
        outline_result = response.finalOutput.artifacts["技能结果"]["esg_report_outline_builder"]
        self.assertIn("GRI", standard_result["reportingBasis"])
        self.assertTrue(matrix_result["disclosureMatrix"])
        self.assertGreaterEqual(kpi_result["indicatorCount"], 3)
        self.assertGreaterEqual(linker_result["linkStats"]["supported"], 1)
        self.assertIn("可再生电力占比", [item["metric"] for item in kpi_result["indicators"]])
        self.assertTrue(all(item["sourceLocation"] for item in kpi_result["indicators"]))
        self.assertIn("# ESG 报告建议大纲", outline_result["outlineMarkdown"])
        self.assertIn("ESG 报告建议大纲", response.finalOutput.revisedDocument or "")
        self.assertTrue(response.finalOutput.evidence)

    def test_execute_esg_report_writer_returns_markdown_and_download(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="请根据这些材料生成一份 800 字 ESG 报告。",
                documents=[
                    {
                        "name": "esg_report_source.txt",
                        "type": "text",
                        "contentText": (
                            "公司 2025 年可再生电力占比提升至 48%，范围二排放同比下降 12%。"
                            "关键岗位培训覆盖率 91%，女性管理者占比 39%。"
                            "Board ESG committee met 4 times. Anti-bribery training coverage reached 100%."
                        ),
                    }
                ],
                manualConfirm=False,
                approved=True,
                agentMode="off",
            )
        )

        self.assertEqual(response.status, "completed")
        self.assertIn("esg_report_writer", response.executedSkills)
        writer_result = response.finalOutput.artifacts["技能结果"]["esg_report_writer"]
        self.assertEqual(writer_result["targetWordCount"], 800)
        self.assertIn("# ESG 报告草稿", writer_result["reportMarkdown"])
        self.assertIn("关键绩效指标", response.finalOutput.revisedDocument or "")
        self.assertTrue(response.finalOutput.downloads)
        self.assertEqual(response.finalOutput.downloads[0].filename, "esg_report_draft.md")

    def test_execute_logs_include_structured_trace_details(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="请总结这份 ESG 材料的重点，并输出一版简短摘要",
                documents=[
                    {
                        "name": "esg_note.txt",
                        "type": "text",
                        "contentText": "公司 2025 年可再生电力占比 48%，培训覆盖率 91%，反舞弊培训覆盖率 100%。",
                    }
                ],
                manualConfirm=True,
                approved=True,
            )
        )

        self.assertEqual(response.status, "completed")
        self.assertGreaterEqual(len(response.logs), 4)

        first_log = response.logs[0]
        self.assertEqual(first_log.kind, "system")
        self.assertEqual(first_log.executor, "workflow_system")
        self.assertEqual(first_log.inputSummary["documentCount"], 1)
        self.assertGreaterEqual(first_log.inputSummary["segmentCount"], 1)

        approval_log = next(log for log in response.logs if log.kind == "checkpoint")
        self.assertEqual(approval_log.executor, "human_review")

        skill_log = next(log for log in response.logs if log.skill == "document_summarizer")
        self.assertEqual(skill_log.kind, "skill")
        self.assertIn(skill_log.executor, {"local_skill", "model_api_agent"})
        self.assertIn("planInputs", skill_log.inputSummary)
        self.assertIn("resultKeys", skill_log.outputSummary)
        self.assertIsNotNone(skill_log.durationMs)

        trace = response.finalOutput.artifacts["执行轨迹"]
        self.assertEqual(len(trace), len(response.logs))
        self.assertEqual(trace[0]["类型"], "系统步骤")
        self.assertIn("输入摘要", trace[0])
        self.assertIn("输出摘要", trace[0])

    def test_execute_requires_confirmation_for_approval_required_skill_even_without_manual_confirm(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="请帮我修订这份说明，整理得更清晰",
                documents=[
                    {
                        "name": "draft.txt",
                        "type": "text",
                        "contentText": "原始草稿内容。",
                    }
                ],
                manualConfirm=False,
                approved=False,
            )
        )

        self.assertEqual(response.status, "awaiting_confirmation")
        self.assertTrue(any(step.checkpoint == "approval" for step in response.plan))

    def test_execute_spreadsheet_confirmation_returns_mapping_preview_before_fill(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="读取原始数据.xlsx 的金额，计算合计和平均值，并填到汇总模板.xlsx。",
                documents=[
                    {
                        "name": "原始数据.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 明细\n项目 | 金额\n打车 | 128\n酒店 | 523.6",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "明细",
                                    "headers": ["项目", "金额"],
                                    "rows": [
                                        {"项目": "打车", "金额": 128},
                                        {"项目": "酒店", "金额": 523.6},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["金额"],
                                }
                            ],
                        },
                    },
                    {
                        "name": "汇总模板.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 汇总表\n指标 | 结果\nsum |\navg |",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "汇总表",
                                    "headers": ["指标", "结果"],
                                    "rows": [
                                        {"指标": "sum", "结果": None},
                                        {"指标": "avg", "结果": None},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": [],
                                }
                            ],
                        },
                    },
                ],
                manualConfirm=True,
                approved=False,
                agentMode="off",
            )
        )

        self.assertEqual(response.status, "awaiting_confirmation")
        self.assertIn("excel_role_classifier", response.executedSkills)
        self.assertIn("calculation_planner", response.executedSkills)
        self.assertIn("spreadsheet_calculator", response.executedSkills)
        self.assertIn("table_mapping_preview", response.executedSkills)
        self.assertNotIn("table_filler", response.executedSkills)
        self.assertNotIn("fill_validator", response.executedSkills)
        self.assertIsNotNone(response.finalOutput)

        preview_result = response.finalOutput.artifacts["技能结果"]["table_mapping_preview"]
        candidate_cells = {item["cell"] for item in preview_result["mappingCandidates"]}
        self.assertEqual(candidate_cells, {"B2", "B3"})

    def test_execute_spreadsheet_workflow_calculates_and_builds_result_table(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="读取这份 Excel，计算金额合计和平均值，再按结果填表输出",
                documents=[
                    {
                        "name": "费用明细.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 费用明细\n项目 | 金额\n打车 | 128\n酒店 | 523.6",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "费用明细",
                                    "headers": ["项目", "金额"],
                                    "rows": [
                                        {"项目": "打车", "金额": 128},
                                        {"项目": "酒店", "金额": 523.6},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["金额"],
                                }
                            ],
                        },
                    }
                ],
                manualConfirm=True,
                approved=True,
            )
        )

        self.assertEqual(response.status, "completed")
        self.assertIn("excel_role_classifier", response.executedSkills)
        self.assertIn("calculation_planner", response.executedSkills)
        self.assertIn("spreadsheet_calculator", response.executedSkills)
        self.assertIn("table_mapping_preview", response.executedSkills)
        self.assertIn("table_filler", response.executedSkills)
        self.assertIn("fill_validator", response.executedSkills)
        spreadsheet_result = response.finalOutput.artifacts["技能结果"]["spreadsheet_calculator"]
        table_result = response.finalOutput.artifacts["技能结果"]["table_filler"]
        validation_result = response.finalOutput.artifacts["技能结果"]["fill_validator"]
        self.assertIn("sum=651.6", spreadsheet_result["summary"])
        self.assertIn("avg=325.8", spreadsheet_result["summary"])
        self.assertTrue(table_result["rows"])
        self.assertIn("| 数值字段 | 指标 | 结果 |", table_result["filledTableMarkdown"])
        self.assertIn("Excel 填写报告", table_result["fillReportMarkdown"])
        self.assertIn("费用明细.xlsx / 费用明细 / 金额", table_result["fillReportMarkdown"])
        self.assertIn("结果表 / E2", table_result["fillReportMarkdown"])
        self.assertTrue(response.finalOutput.downloads)
        self.assertTrue(response.finalOutput.downloads[0].filename.endswith(".xlsx"))
        fill_report_download = next(
            item for item in response.finalOutput.downloads if item.filename.endswith("_fill_report.md")
        )
        fill_report_text = base64.b64decode(fill_report_download.contentBase64).decode("utf-8")
        self.assertIn("Excel 填写报告", fill_report_text)
        self.assertIn("651.6", fill_report_text)
        self.assertIn("结果表 / E2", fill_report_text)

        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(response.finalOutput.downloads[0].contentBase64))
        )
        self.assertIn("结果表", workbook.sheetnames)
        result_sheet = workbook["结果表"]
        self.assertEqual(result_sheet["A1"].value, "文档")
        self.assertEqual(result_sheet["D2"].value, "sum")
        self.assertEqual(result_sheet["E2"].value, 651.6)
        self.assertEqual(validation_result["validationStats"]["failed"], 0)
        self.assertTrue(response.finalOutput.nextActions)

    def test_fill_validator_handles_unreadable_export_content(self) -> None:
        result = FillValidatorSkill().execute(
            SkillExecutionContext(
                task="校验损坏的导出文件",
                intention=IntentionAnalysis(
                    primaryGoal="校验损坏的导出文件",
                    intentType="review",
                    detectedIntentTypes=["review"],
                    confidence=0.8,
                ),
                documents=[],
                previous_results={
                    "table_filler": {
                        "exportFiles": [{"contentBase64": "not-base64"}],
                        "fillAudit": [
                            {
                                "field": "结果",
                                "status": "written",
                                "sheet": "汇总表",
                                "cell": "B2",
                                "value": 100,
                            }
                        ],
                    }
                },
            )
        )

        self.assertEqual(result["validationStats"]["checked"], 0)
        self.assertTrue(result["validationStats"]["readError"])
        self.assertIn("无法读取", result["summary"])

    def test_table_filler_prefers_existing_template_sheet_with_metric_and_result_columns(self) -> None:
        response = workflow_agent_service.execute(
            request=WorkflowExecuteRequest(
                task="读取这份 Excel，计算金额合计和平均值，并填到汇总表。",
                documents=[
                    {
                        "name": "费用模板.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 明细\n项目 | 金额\n打车 | 128\n酒店 | 523.6\n\n[工作表] 汇总表\n指标 | 结果",
                        "structuredData": {
                            "sheetCount": 2,
                            "sheets": [
                                {
                                    "title": "明细",
                                    "headers": ["项目", "金额"],
                                    "rows": [
                                        {"项目": "打车", "金额": 128},
                                        {"项目": "酒店", "金额": 523.6},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["金额"],
                                },
                                {
                                    "title": "汇总表",
                                    "headers": ["指标", "结果"],
                                    "rows": [
                                        {"指标": "sum", "结果": None},
                                        {"指标": "avg", "结果": None},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": [],
                                },
                            ],
                        },
                    }
                ],
                manualConfirm=False,
                approved=True,
            )
        )

        self.assertEqual(response.status, "completed")
        self.assertTrue(response.finalOutput.downloads)
        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(response.finalOutput.downloads[0].contentBase64))
        )
        self.assertIn("汇总表", workbook.sheetnames)
        summary_sheet = workbook["汇总表"]
        self.assertEqual(summary_sheet["A2"].value, "sum")
        self.assertEqual(summary_sheet["B2"].value, 651.6)
        self.assertEqual(summary_sheet["A3"].value, "avg")
        self.assertEqual(summary_sheet["B3"].value, 325.8)


class WorkflowApiTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.client = TestClient(app)

    def setUp(self) -> None:
        self.table_fill_review_patcher = patch(
            "app.services.workflow.table_fill.crosscheck.TableFillMappingCrossCheckService.review_mapping_plan",
            return_value={
                "enabled": False,
                "status": "skipped",
                "approved": True,
                "riskLevel": "unknown",
                "blockWrite": False,
                "issues": [],
                "suggestions": [],
                "candidateReviews": [],
            },
        )
        self.table_fill_review_patcher.start()
        self.addCleanup(self.table_fill_review_patcher.stop)

    def _build_docx_bytes(self, lines: list[str]) -> bytes:
        buffer = BytesIO()
        document = Document()
        for line in lines:
            document.add_paragraph(line)
        document.save(buffer)
        return buffer.getvalue()

    def _build_xlsx_bytes(
        self,
        *,
        headers: list[str] | None = None,
        data_rows: list[list[object]] | None = None,
    ) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "报销明细"
        sheet.append(headers or ["项目", "金额"])
        for row in data_rows or [["打车", 128], ["酒店", 523.6]]:
            sheet.append(row)
        workbook.save(buffer)
        return buffer.getvalue()

    def _build_direct_transfer_target_xlsx_bytes(self) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "目标填报"
        sheet.append(["项目", "金额"])
        workbook.save(buffer)
        return buffer.getvalue()

    def _build_template_xlsx_bytes(self) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()

        detail_sheet = workbook.active
        detail_sheet.title = "明细"
        detail_sheet.append(["项目", "金额"])
        detail_sheet.append(["打车", 128])
        detail_sheet.append(["酒店", 523.6])

        summary_sheet = workbook.create_sheet("汇总表")
        summary_sheet["A1"] = "费用汇总"
        summary_sheet["A2"] = "指标"
        summary_sheet["B2"] = "结果"
        summary_sheet["A3"] = "sum"
        summary_sheet["A4"] = "avg"

        workbook.save(buffer)
        return buffer.getvalue()

    def _build_summary_template_xlsx_bytes(self) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()

        summary_sheet = workbook.active
        summary_sheet.title = "汇总表"
        summary_sheet["A1"] = "费用汇总"
        summary_sheet["A2"] = "指标"
        summary_sheet["B2"] = "结果"
        summary_sheet["A3"] = "sum"
        summary_sheet["A4"] = "avg"

        workbook.save(buffer)
        return buffer.getvalue()

    def _build_label_value_template_xlsx_bytes(self, *, existing_sum: object | None = None) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()

        summary_sheet = workbook.active
        summary_sheet.title = "汇总页"
        summary_sheet["A1"] = "费用汇总"
        summary_sheet["A2"] = "合计"
        summary_sheet["B2"] = existing_sum
        summary_sheet["A3"] = "平均值"
        summary_sheet["B3"] = None

        workbook.save(buffer)
        return buffer.getvalue()

    def _build_contextual_label_template_xlsx_bytes(self) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()

        summary_sheet = workbook.active
        summary_sheet.title = "填报页"
        summary_sheet["A1"] = "业务汇总"
        summary_sheet["A2"] = "金额合计"
        summary_sheet["A3"] = "人数合计"

        workbook.save(buffer)
        return buffer.getvalue()

    def _build_grouped_label_template_xlsx_bytes(self) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()

        summary_sheet = workbook.active
        summary_sheet.title = "部门汇总"
        summary_sheet["A1"] = "部门费用汇总"
        summary_sheet["A2"] = "销售部合计"
        summary_sheet["A3"] = "研发部合计"

        workbook.save(buffer)
        return buffer.getvalue()

    def _build_many_group_source_xlsx_bytes(self, *, count: int) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "项目明细"
        sheet.append(["项目", "金额"])
        for index in range(1, count + 1):
            sheet.append([f"项目{index:03d}", index])
        workbook.save(buffer)
        return buffer.getvalue()

    def _build_many_group_label_template_xlsx_bytes(self, *, count: int) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()
        summary_sheet = workbook.active
        summary_sheet.title = "批量填报"
        summary_sheet["A1"] = "项目金额填报"
        for index in range(1, count + 1):
            summary_sheet.cell(row=index + 1, column=1).value = f"项目{index:03d}合计"
        workbook.save(buffer)
        return buffer.getvalue()

    def _build_header_offset_xlsx_bytes(self) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "排放汇总"
        sheet["A1"] = "2025 ESG 排放汇总表"
        sheet.merge_cells("A1:C1")
        sheet["A3"] = "月份"
        sheet["B3"] = "排放量"
        sheet["C3"] = "单位"
        sheet.append([None, None, None])
        sheet.append(["1月", 128, "吨"])
        sheet.append(["2月", 132, "吨"])
        workbook.save(buffer)
        return buffer.getvalue()

    def _build_merged_label_template_xlsx_bytes(self, *, with_blocked_right_cells: bool = False, existing_sum: object | None = None) -> bytes:
        buffer = BytesIO()
        workbook = Workbook()
        summary_sheet = workbook.active
        summary_sheet.title = "汇总页"
        summary_sheet["A1"] = "费用汇总"
        summary_sheet.merge_cells("A1:B1")
        summary_sheet["A2"] = "合计"
        summary_sheet["B2"] = "平均值" if with_blocked_right_cells else existing_sum
        summary_sheet["C2"] = "同比" if with_blocked_right_cells else None
        summary_sheet["D2"] = "环比" if with_blocked_right_cells else None
        summary_sheet["A3"] = "平均值"
        summary_sheet["B3"] = "环比" if with_blocked_right_cells else None
        workbook.save(buffer)
        return buffer.getvalue()

    def _build_pptx_bytes(self, slides: list[str]) -> bytes:
        from pptx import Presentation

        buffer = BytesIO()
        presentation = Presentation()
        for slide_text in slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "ESG 幻灯片"
            slide.placeholders[1].text = slide_text
        presentation.save(buffer)
        return buffer.getvalue()

    def _wait_for_job_terminal_state(self, job_id: str) -> dict:
        for _ in range(80):
            payload = self.client.get(f"/api/workflow/jobs/{job_id}").json()
            if payload["status"] in {"completed", "blocked", "failed", "awaiting_confirmation"}:
                return payload
            time.sleep(0.05)
        self.fail(f"Job {job_id} did not reach a terminal state in time.")

    def test_workflow_skills_endpoint(self) -> None:
        response = self.client.get("/api/workflow/skills")

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertGreaterEqual(payload["total"], 6)

    def test_upload_endpoint_reads_text_file(self) -> None:
        response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    ("notes.txt", "第一行材料说明\n第二行费用详情".encode("utf-8"), "text/plain"),
                )
            ],
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["total"], 1)
        self.assertEqual(payload["documents"][0]["type"], "text")
        self.assertEqual(payload["documents"][0]["parser"], "text_parser")
        self.assertIn("第一行材料说明", payload["documents"][0]["contentText"])
        self.assertGreaterEqual(len(payload["documents"][0]["segments"]), 1)
        self.assertEqual(payload["documents"][0]["segments"][0]["kind"], "paragraph")

    def test_upload_endpoint_reads_docx_and_xlsx(self) -> None:
        response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "说明.docx",
                        self._build_docx_bytes(["这是第一段。", "这是第二段。"]),
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    ),
                ),
                (
                    "files",
                    (
                        "明细.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["total"], 2)
        by_name = {item["name"]: item for item in payload["documents"]}
        self.assertEqual(by_name["说明.docx"]["parser"], "docx_parser")
        self.assertIn("这是第一段。", by_name["说明.docx"]["contentText"])
        self.assertGreaterEqual(len(by_name["说明.docx"]["segments"]), 2)
        self.assertEqual(by_name["明细.xlsx"]["parser"], "excel_parser")
        self.assertIn("报销明细", by_name["明细.xlsx"]["contentText"])
        self.assertGreaterEqual(len(by_name["明细.xlsx"]["segments"]), 2)
        self.assertEqual(by_name["明细.xlsx"]["segments"][0]["kind"], "sheet_summary")
        self.assertEqual(by_name["明细.xlsx"]["structuredData"]["sheetCount"], 1)
        self.assertEqual(by_name["明细.xlsx"]["structuredData"]["sheets"][0]["headers"], ["项目", "金额"])
        self.assertEqual(by_name["明细.xlsx"]["structuredData"]["sheets"][0]["numericColumns"], ["金额"])
        self.assertEqual(by_name["明细.xlsx"]["structuredData"]["sheets"][0]["rows"][0]["项目"], "打车")
        self.assertTrue(by_name["明细.xlsx"]["structuredData"]["workbookUploadToken"])

    def test_upload_endpoint_extracts_excel_layout_metadata(self) -> None:
        response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "排放汇总.xlsx",
                        self._build_header_offset_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                )
            ],
        )

        self.assertEqual(response.status_code, 200)
        document = response.json()["documents"][0]
        sheet = document["structuredData"]["sheets"][0]
        self.assertEqual(sheet["headerRowIndex"], 3)
        self.assertIn(sheet["headerConfidence"], {"medium", "high"})
        self.assertEqual(sheet["sheetRole"], "source_like")
        self.assertTrue(sheet["mergedRanges"])
        self.assertIn("识别表头位于第 3 行", sheet["layoutSummary"])

    def test_upload_endpoint_extracts_label_anchors_and_empty_value_zones(self) -> None:
        response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "汇总模板.xlsx",
                        self._build_merged_label_template_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                )
            ],
        )

        self.assertEqual(response.status_code, 200)
        document = response.json()["documents"][0]
        sheet = document["structuredData"]["sheets"][0]
        anchor_metrics = {item["metric"] for item in sheet["labelAnchors"]}
        zone_cells = {item["cell"] for item in sheet["emptyValueZones"]}
        self.assertTrue({"sum", "avg"}.issubset(anchor_metrics))
        self.assertIn("A1:B1", {item["range"] for item in sheet["mergedRanges"]})
        self.assertTrue({"B2", "B3"} <= zone_cells)

    def test_upload_endpoint_reads_pptx_file(self) -> None:
        response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "汇报材料.pptx",
                        self._build_pptx_bytes(["可再生电力占比 48%。", "反舞弊培训覆盖率 100%。"]),
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    ),
                )
            ],
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        document = payload["documents"][0]
        self.assertEqual(document["type"], "presentation")
        self.assertEqual(document["parser"], "pptx_parser")
        self.assertGreaterEqual(len(document["segments"]), 2)
        self.assertIn("可再生电力占比 48%", document["contentText"])

    def test_upload_endpoint_reads_zip_bundle(self) -> None:
        buffer = BytesIO()
        with zipfile.ZipFile(buffer, "w") as archive:
            archive.writestr("climate.txt", "可再生电力占比 48%。")
            archive.writestr("governance.txt", "Board ESG committee met 4 times.")

        response = self.client.post(
            "/api/workflow/uploads",
            files=[("files", ("bundle.zip", buffer.getvalue(), "application/zip"))],
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        document = payload["documents"][0]
        self.assertEqual(document["type"], "archive")
        self.assertEqual(document["parser"], "zip_parser")
        self.assertGreaterEqual(len(document["segments"]), 3)
        self.assertIn("可再生电力占比 48%", document["contentText"])

    def test_execute_api_keeps_original_template_when_multiple_excels_are_uploaded(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "费用模板.xlsx",
                        self._build_template_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "其他台账.xlsx",
                        self._build_xlsx_bytes(
                            headers=["月份", "人数"],
                            data_rows=[["一月", 12], ["二月", 15]],
                        ),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取费用模板.xlsx，计算金额合计和平均值，并填到汇总表。",
                    "documents": uploaded_documents,
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        self.assertEqual(payload["status"], "completed")
        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(payload["finalOutput"]["downloads"][0]["contentBase64"]))
        )
        self.assertIn("汇总表", workbook.sheetnames)
        summary_sheet = workbook["汇总表"]
        self.assertEqual(summary_sheet["A1"].value, "费用汇总")
        self.assertEqual(summary_sheet["A3"].value, "sum")
        self.assertEqual(summary_sheet["B3"].value, 651.6)
        self.assertEqual(summary_sheet["A4"].value, "avg")
        self.assertEqual(summary_sheet["B4"].value, 325.8)

    def test_upload_endpoint_runs_ocr_for_image_files(self) -> None:
        mocked_ocr = OCRExtraction(
            text="出租车发票\n金额 128.00",
            notes=["已使用 阿里云百炼 OCR 提取图片正文。"],
            segments=[
                DocumentSegment(
                    segmentId="ocr_1",
                    kind="ocr",
                    label="OCR 片段 1",
                    text="出租车发票\n金额 128.00",
                )
            ],
        )

        with patch(
            "app.services.workflow.uploads.parsers.workflow_ocr_service.extract_image_text",
            return_value=mocked_ocr,
        ):
            response = self.client.post(
                "/api/workflow/uploads",
                files=[("files", ("receipt.png", b"fake-image-bytes", "image/png"))],
            )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        document = payload["documents"][0]
        self.assertEqual(document["parser"], "image_ocr_parser")
        self.assertEqual(document["ocrText"], "出租车发票\n金额 128.00")
        self.assertIn("金额 128.00", document["contentText"])
        self.assertEqual(document["segments"][0]["kind"], "ocr")

    def test_upload_endpoint_keeps_long_text_without_truncation(self) -> None:
        long_text = "ESG 披露数据。" * 3000

        response = self.client.post(
            "/api/workflow/uploads",
            files=[("files", ("long_note.txt", long_text.encode("utf-8"), "text/plain"))],
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        document = payload["documents"][0]
        self.assertEqual(document["parser"], "text_parser")
        self.assertEqual(document["contentText"], long_text)
        self.assertNotIn("内容过长，已截断", document["contentText"])

    def test_upload_endpoint_keeps_full_spreadsheet_rows_and_columns(self) -> None:
        headers = [f"列{index}" for index in range(1, 14)]
        data_rows = [[f"值{row}_{column}" for column in range(1, 14)] for row in range(1, 61)]

        response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "full.xlsx",
                        self._build_xlsx_bytes(headers=headers, data_rows=data_rows),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                )
            ],
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        document = payload["documents"][0]
        self.assertEqual(document["structuredData"]["sheets"][0]["rowCount"], 60)
        self.assertEqual(len(document["structuredData"]["sheets"][0]["headers"]), 13)
        self.assertEqual(document["structuredData"]["sheets"][0]["headers"][-1], "列13")
        self.assertEqual(document["structuredData"]["sheets"][0]["rows"][-1]["列13"], "值60_13")

    def test_execute_api_can_calculate_from_uploaded_csv(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "费用明细.csv",
                        "项目,金额\n打车,128\n酒店,523.6".encode("utf-8"),
                        "text/csv",
                    ),
                )
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_document = upload_response.json()["documents"][0]
        self.assertEqual(uploaded_document["parser"], "csv_parser")
        self.assertEqual(uploaded_document["structuredData"]["sheetCount"], 1)
        self.assertEqual(uploaded_document["structuredData"]["sheets"][0]["headers"], ["项目", "金额"])
        self.assertEqual(uploaded_document["structuredData"]["sheets"][0]["numericColumns"], ["金额"])
        self.assertNotIn("workbookUploadToken", uploaded_document["structuredData"])

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取这个 CSV，计算金额合计和平均值，再按结果填表输出",
                    "documents": [uploaded_document],
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        self.assertEqual(payload["status"], "completed")
        self.assertIn("spreadsheet_calculator", payload["executedSkills"])
        self.assertIn("table_filler", payload["executedSkills"])
        self.assertIn(
            "sum=651.6",
            payload["finalOutput"]["artifacts"]["技能结果"]["spreadsheet_calculator"]["summary"],
        )
        self.assertTrue(payload["finalOutput"]["downloads"])

    def test_session_api_persists_uploaded_documents_plan_and_execution(self) -> None:
        session_response = self.client.post("/api/workflow/sessions")

        self.assertEqual(session_response.status_code, 200)
        session_payload = session_response.json()
        session_id = session_payload["sessionId"]
        self.assertTrue(session_id)
        self.assertEqual(session_payload["documents"], [])

        upload_response = self.client.post(
            "/api/workflow/uploads",
            data={"sessionId": session_id},
            files=[
                (
                    "files",
                    ("notes.txt", "董事会 ESG 委员会召开 4 次。".encode("utf-8"), "text/plain"),
                )
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        upload_payload = upload_response.json()
        self.assertEqual(upload_payload["sessionId"], session_id)
        self.assertEqual(len(upload_payload["documents"]), 1)
        self.assertEqual(len(upload_payload["mergedDocuments"]), 1)

        session_state_response = self.client.get(f"/api/workflow/sessions/{session_id}")
        self.assertEqual(session_state_response.status_code, 200)
        session_state = session_state_response.json()
        self.assertEqual(len(session_state["documents"]), 1)
        self.assertEqual(session_state["documents"][0]["name"], "notes.txt")

        uploaded_document = upload_payload["mergedDocuments"][0]
        plan_response = self.client.post(
            "/api/workflow/plan",
            json={
                "sessionId": session_id,
                "task": "请总结这份治理说明",
                "documents": [uploaded_document],
                "manualConfirm": False,
                "agentMode": "off",
                "localFallbackEnabled": True,
                "preferredSkills": [],
                "context": {},
            },
        )

        self.assertEqual(plan_response.status_code, 200)
        plan_payload = plan_response.json()
        self.assertEqual(plan_payload["sessionId"], session_id)

        session_state_response = self.client.get(f"/api/workflow/sessions/{session_id}")
        session_state = session_state_response.json()
        self.assertEqual(session_state["task"], "请总结这份治理说明")
        self.assertTrue(session_state["localFallbackEnabled"])
        self.assertEqual(session_state["latestPlanResponse"]["sessionId"], session_id)

        execute_response = self.client.post(
            "/api/workflow/execute",
            json={
                "sessionId": session_id,
                "task": "请总结这份治理说明",
                "documents": [uploaded_document],
                "manualConfirm": False,
                "approved": True,
                "agentMode": "off",
                "localFallbackEnabled": True,
                "preferredSkills": [],
                "context": {},
            },
        )

        self.assertEqual(execute_response.status_code, 200)
        execute_payload = execute_response.json()
        self.assertEqual(execute_payload["sessionId"], session_id)
        self.assertEqual(execute_payload["status"], "completed")

        session_state_response = self.client.get(f"/api/workflow/sessions/{session_id}")
        session_state = session_state_response.json()
        self.assertEqual(session_state["latestExecutionResponse"]["sessionId"], session_id)
        self.assertEqual(session_state["latestExecutionResponse"]["status"], "completed")

    def test_plan_api_returns_error_when_fallback_disabled_and_agent_has_no_result(self) -> None:
        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None):
            response = self.client.post(
                "/api/workflow/plan",
                json={
                    "task": "请帮我总结这份说明",
                    "documents": [
                        {
                            "name": "notes.txt",
                            "type": "text",
                            "contentText": "这是需要整理的说明文稿。",
                        }
                    ],
                    "manualConfirm": False,
                    "agentMode": "on",
                    "localFallbackEnabled": False,
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(response.status_code, 503)
        payload = response.json()
        self.assertIn("本地回退", payload["detail"])

    def test_delete_session_endpoint_clears_session_state(self) -> None:
        session_response = self.client.post("/api/workflow/sessions")
        session_id = session_response.json()["sessionId"]

        delete_response = self.client.delete(f"/api/workflow/sessions/{session_id}")
        self.assertEqual(delete_response.status_code, 200)
        self.assertEqual(delete_response.json()["sessionId"], session_id)

        missing_response = self.client.get(f"/api/workflow/sessions/{session_id}")
        self.assertEqual(missing_response.status_code, 404)

    def test_execute_api_fills_uploaded_template_workbook_and_preserves_title_row(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "费用模板.xlsx",
                        self._build_template_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                )
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_document = upload_response.json()["documents"][0]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取这份 Excel，计算金额合计和平均值，并填到汇总表。",
                    "documents": [uploaded_document],
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        self.assertEqual(payload["status"], "completed")
        downloads = payload["finalOutput"]["downloads"]
        self.assertTrue(downloads)

        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(downloads[0]["contentBase64"]))
        )
        self.assertIn("汇总表", workbook.sheetnames)
        summary_sheet = workbook["汇总表"]
        self.assertEqual(summary_sheet["A1"].value, "费用汇总")
        self.assertEqual(summary_sheet["A2"].value, "指标")
        self.assertEqual(summary_sheet["B2"].value, "结果")
        self.assertEqual(summary_sheet["A3"].value, "sum")
        self.assertEqual(summary_sheet["B3"].value, 651.6)
        self.assertEqual(summary_sheet["A4"].value, "avg")
        self.assertEqual(summary_sheet["B4"].value, 325.8)

    def test_execute_api_fills_separate_template_workbook_from_source_excel(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "原始数据.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "回填模板.xlsx",
                        self._build_summary_template_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取原始数据.xlsx 中的金额，计算合计和平均值，并填到回填模板.xlsx 的汇总表。",
                    "documents": uploaded_documents,
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        self.assertEqual(payload["status"], "completed")
        downloads = payload["finalOutput"]["downloads"]
        self.assertTrue(downloads)
        self.assertEqual(downloads[0]["filename"], "回填模板_filled.xlsx")

        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(downloads[0]["contentBase64"]))
        )
        self.assertEqual(workbook.sheetnames, ["汇总表"])
        summary_sheet = workbook["汇总表"]
        self.assertEqual(summary_sheet["A1"].value, "费用汇总")
        self.assertEqual(summary_sheet["A3"].value, "sum")
        self.assertEqual(summary_sheet["B3"].value, 651.6)
        self.assertEqual(summary_sheet["A4"].value, "avg")
        self.assertEqual(summary_sheet["B4"].value, 325.8)

    def test_execute_api_transfers_source_table_rows_into_target_table(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "数据源表格.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "目标表格.xlsx",
                        self._build_direct_transfer_target_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None), \
             patch(
                 "app.services.workflow.table_fill.crosscheck.TableDataTransferCrossCheckService.review_transfer_plan",
                 return_value={
                     "enabled": False,
                     "provider": "openai",
                     "providerLabel": "OpenAI",
                     "model": "gpt-5.4-mini",
                     "status": "skipped",
                     "approved": True,
                     "riskLevel": "unknown",
                     "blockWrite": False,
                     "issues": [],
                     "suggestions": [],
                 },
             ):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "将数据源表格.xlsx 里的数据自动写入目标表格.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        self.assertEqual(payload["status"], "completed")
        self.assertIn("excel_role_classifier", payload["executedSkills"])
        self.assertIn("table_data_transfer", payload["executedSkills"])
        self.assertNotIn("spreadsheet_calculator", payload["executedSkills"])
        self.assertNotIn("table_filler", payload["executedSkills"])

        downloads = payload["finalOutput"]["downloads"]
        self.assertTrue(downloads)
        self.assertEqual(downloads[0]["filename"], "目标表格_filled.xlsx")

        workbook = load_workbook(filename=BytesIO(base64.b64decode(downloads[0]["contentBase64"])))
        target_sheet = workbook["目标填报"]
        self.assertEqual(target_sheet["A1"].value, "项目")
        self.assertEqual(target_sheet["B1"].value, "金额")
        self.assertEqual(target_sheet["A2"].value, "打车")
        self.assertEqual(target_sheet["B2"].value, 128)
        self.assertEqual(target_sheet["A3"].value, "酒店")
        self.assertEqual(target_sheet["B3"].value, 523.6)

        transfer_result = payload["finalOutput"]["artifacts"]["技能结果"]["table_data_transfer"]
        self.assertEqual(transfer_result["targetSheet"], "目标填报")
        self.assertFalse(transfer_result["crossCheck"]["enabled"])
        self.assertEqual(transfer_result["crossCheck"]["status"], "skipped")
        self.assertEqual(transfer_result["transferStats"]["rowsTransferred"], 2)
        self.assertEqual(transfer_result["transferStats"]["written"], 4)
        self.assertEqual(
            [item["sourceHeader"] for item in transfer_result["columnMappings"]],
            ["项目", "金额"],
        )

    def test_execute_api_fills_label_value_template_without_headers(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "原始数据.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "空白汇总模板.xlsx",
                        self._build_label_value_template_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取原始数据.xlsx 的金额，计算合计和平均值，并填到空白汇总模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        downloads = payload["finalOutput"]["downloads"]
        self.assertTrue(downloads)

        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(downloads[0]["contentBase64"]))
        )
        summary_sheet = workbook["汇总页"]
        self.assertEqual(summary_sheet["B2"].value, 651.6)
        self.assertEqual(summary_sheet["B3"].value, 325.8)

        table_result = payload["finalOutput"]["artifacts"]["技能结果"]["table_filler"]
        self.assertEqual(table_result["fillStats"]["mode"], "label_value")
        audit_by_cell = {item["cell"]: item for item in table_result["fillAudit"]}
        self.assertEqual(audit_by_cell["B2"]["status"], "written")
        self.assertEqual(audit_by_cell["B3"]["status"], "written")

    def test_execute_api_fills_contextual_labels_when_metric_names_repeat(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "原始数据.xlsx",
                        self._build_xlsx_bytes(
                            headers=["项目", "金额", "人数"],
                            data_rows=[["门店A", 128, 3], ["门店B", 523.6, 7]],
                        ),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "上下文标签模板.xlsx",
                        self._build_contextual_label_template_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取原始数据.xlsx，计算金额和人数合计，并填到上下文标签模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(payload["finalOutput"]["downloads"][0]["contentBase64"]))
        )
        summary_sheet = workbook["填报页"]
        self.assertIsNone(summary_sheet["B1"].value)
        self.assertEqual(summary_sheet["B2"].value, 651.6)
        self.assertEqual(summary_sheet["B3"].value, 10)

        table_result = payload["finalOutput"]["artifacts"]["技能结果"]["table_filler"]
        audit_by_mapping_id = {item["mappingId"]: item for item in table_result["fillAudit"] if item.get("mappingId")}
        self.assertEqual(audit_by_mapping_id["map_1"]["cell"], "B2")
        self.assertEqual(audit_by_mapping_id["map_2"]["cell"], "B3")

    def test_execute_api_fills_grouped_result_labels_by_group_value(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "部门费用.xlsx",
                        self._build_xlsx_bytes(
                            headers=["部门", "金额"],
                            data_rows=[["销售部", 120], ["研发部", 80], ["销售部", 30]],
                        ),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "分组标签模板.xlsx",
                        self._build_grouped_label_template_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取部门费用.xlsx，按部门计算金额合计，并填到分组标签模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(payload["finalOutput"]["downloads"][0]["contentBase64"]))
        )
        summary_sheet = workbook["部门汇总"]
        self.assertIsNone(summary_sheet["B1"].value)
        self.assertEqual(summary_sheet["B2"].value, 150)
        self.assertEqual(summary_sheet["B3"].value, 80)

    def test_execute_api_fills_many_label_value_template_rows(self) -> None:
        item_count = 120
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "大量项目明细.xlsx",
                        self._build_many_group_source_xlsx_bytes(count=item_count),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "大量填报模板.xlsx",
                        self._build_many_group_label_template_xlsx_bytes(count=item_count),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取大量项目明细.xlsx，按项目计算金额合计，并填到大量填报模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(payload["finalOutput"]["downloads"][0]["contentBase64"]))
        )
        summary_sheet = workbook["批量填报"]
        self.assertEqual(summary_sheet["B2"].value, 1)
        self.assertEqual(summary_sheet["B41"].value, 40)
        self.assertEqual(summary_sheet["B121"].value, 120)

        table_result = payload["finalOutput"]["artifacts"]["技能结果"]["table_filler"]
        self.assertEqual(table_result["fillStats"]["mode"], "label_value")
        self.assertEqual(table_result["fillStats"]["written"], item_count)
        written_cells = {
            item["cell"]
            for item in table_result["fillAudit"]
            if item.get("status") == "written"
        }
        self.assertIn("B121", written_cells)

    def test_execute_api_preserves_existing_template_values_and_records_audit(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "原始数据.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "冲突模板.xlsx",
                        self._build_label_value_template_xlsx_bytes(existing_sum=600),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            execute_response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取原始数据.xlsx 的金额，计算合计和平均值，并填到冲突模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": False,
                    "approved": True,
                    "agentMode": "off",
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(execute_response.status_code, 200)
        payload = execute_response.json()
        downloads = payload["finalOutput"]["downloads"]
        self.assertTrue(downloads)

        workbook = load_workbook(
            filename=BytesIO(base64.b64decode(downloads[0]["contentBase64"]))
        )
        summary_sheet = workbook["汇总页"]
        self.assertEqual(summary_sheet["B2"].value, 600)
        self.assertEqual(summary_sheet["B3"].value, 325.8)

        table_result = payload["finalOutput"]["artifacts"]["技能结果"]["table_filler"]
        self.assertEqual(table_result["fillStats"]["preservedExisting"], 1)
        audit_by_cell = {item["cell"]: item for item in table_result["fillAudit"]}
        self.assertEqual(audit_by_cell["B2"]["status"], "preserved_existing")
        self.assertEqual(audit_by_cell["B3"]["status"], "written")

    def test_job_api_confirmation_preview_allows_manual_mapping_override(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "原始数据.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "空白汇总模板.xlsx",
                        self._build_label_value_template_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        calculator_skill = workflow_agent_service._skill_registry.require("spreadsheet_calculator")
        preview_skill = workflow_agent_service._skill_registry.require("table_mapping_preview")

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None), \
             patch.object(calculator_skill, "execute", wraps=calculator_skill.execute) as calculator_execute, \
             patch.object(preview_skill, "execute", wraps=preview_skill.execute) as preview_execute:
            create_response = self.client.post(
                "/api/workflow/jobs",
                json={
                    "task": "读取原始数据.xlsx 的金额，计算合计和平均值，并填到空白汇总模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": True,
                    "approved": False,
                    "agentMode": "off",
                    "localFallbackEnabled": True,
                    "preferredSkills": [],
                    "context": {},
                    "planOverrides": {},
                },
            )

            self.assertEqual(create_response.status_code, 200)
            job_id = create_response.json()["jobId"]
            awaiting = self._wait_for_job_terminal_state(job_id)
            self.assertEqual(awaiting["status"], "awaiting_confirmation")
            self.assertIn("table_mapping_preview", awaiting["executedSkills"])
            self.assertNotIn("table_filler", awaiting["executedSkills"])
            self.assertEqual(calculator_execute.call_count, 1)
            self.assertEqual(preview_execute.call_count, 1)

            preview_result = awaiting["finalOutput"]["artifacts"]["技能结果"]["table_mapping_preview"]
            preview_by_id = {item["mappingId"]: item for item in preview_result["mappingCandidates"]}
            self.assertEqual(preview_by_id["map_1"]["cell"], "B2")
            self.assertEqual(preview_by_id["map_2"]["cell"], "B3")

            table_filler_step = next(step for step in awaiting["plan"] if step.get("skill") == "table_filler")
            approve_response = self.client.post(
                f"/api/workflow/jobs/{job_id}/approve",
                json={
                    "approved": True,
                    "planOverrides": {
                        "disabledStepIds": [],
                        "stepInputOverrides": {
                            table_filler_step["stepId"]: {
                                "manualMappings": [
                                    {"mappingId": "map_1", "sheet": "汇总页", "cell": "C2"},
                                    {"mappingId": "map_2", "sheet": "汇总页", "cell": "C3"},
                                ]
                            }
                        },
                    },
                },
            )
            self.assertEqual(approve_response.status_code, 200)

            completed = self._wait_for_job_terminal_state(job_id)
            self.assertEqual(completed["status"], "completed")
            self.assertEqual(calculator_execute.call_count, 1)
            self.assertEqual(preview_execute.call_count, 1)
            downloads = completed["finalOutput"]["downloads"]
            self.assertTrue(downloads)

            artifact_response = self.client.get(downloads[0]["downloadUrl"])
            self.assertEqual(artifact_response.status_code, 200)
            workbook = load_workbook(filename=BytesIO(artifact_response.content))
            summary_sheet = workbook["汇总页"]
            self.assertIsNone(summary_sheet["B2"].value)
            self.assertIsNone(summary_sheet["B3"].value)
            self.assertEqual(summary_sheet["C2"].value, 651.6)
            self.assertEqual(summary_sheet["C3"].value, 325.8)

            table_result = completed["finalOutput"]["artifacts"]["技能结果"]["table_filler"]
            audit_by_cell = {item["cell"]: item for item in table_result["fillAudit"] if item.get("cell")}
            self.assertEqual(audit_by_cell["C2"]["status"], "written")
            self.assertEqual(audit_by_cell["C3"]["status"], "written")
            self.assertEqual(table_result["fillStats"]["mode"], "manual_mapping")

    def test_confirmation_preview_includes_candidate_explanations_and_alternatives(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "原始数据.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "冲突模板.xlsx",
                        self._build_label_value_template_xlsx_bytes(existing_sum=600),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            response = self.client.post(
                "/api/workflow/execute",
                json={
                    "task": "读取原始数据.xlsx 的金额，计算合计和平均值，并填到冲突模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": True,
                    "approved": False,
                    "agentMode": "off",
                    "localFallbackEnabled": True,
                    "preferredSkills": [],
                    "context": {},
                },
            )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        preview_result = payload["finalOutput"]["artifacts"]["技能结果"]["table_mapping_preview"]
        sum_candidate = next(item for item in preview_result["mappingCandidates"] if item["metric"] == "sum")
        self.assertEqual(sum_candidate["topCandidate"]["cell"], "B2")
        self.assertTrue(sum_candidate["alternativeCandidates"])
        self.assertEqual(sum_candidate["alternativeCandidates"][0]["cell"], "C2")
        self.assertIsInstance(sum_candidate["score"], float)
        self.assertTrue(sum_candidate["reasons"])
        self.assertIn(sum_candidate["riskLevel"], {"low", "medium", "high"})

    def test_job_api_blocks_low_confidence_candidate_without_explicit_confirmation(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "原始数据.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "高风险模板.xlsx",
                        self._build_merged_label_template_xlsx_bytes(with_blocked_right_cells=True),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            create_response = self.client.post(
                "/api/workflow/jobs",
                json={
                    "task": "读取原始数据.xlsx 的金额，计算合计，并填到高风险模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": True,
                    "approved": False,
                    "agentMode": "off",
                    "localFallbackEnabled": True,
                    "preferredSkills": [],
                    "context": {},
                    "planOverrides": {},
                },
            )

            self.assertEqual(create_response.status_code, 200)
            job_id = create_response.json()["jobId"]
            awaiting = self._wait_for_job_terminal_state(job_id)
            preview_result = awaiting["finalOutput"]["artifacts"]["技能结果"]["table_mapping_preview"]
            self.assertTrue(any(item["requiresConfirmation"] for item in preview_result["mappingCandidates"]))

            table_filler_step = next(step for step in awaiting["plan"] if step.get("skill") == "table_filler")
            approve_response = self.client.post(
                f"/api/workflow/jobs/{job_id}/approve",
                json={
                    "approved": True,
                    "planOverrides": {
                        "disabledStepIds": [],
                        "stepInputOverrides": {
                            table_filler_step["stepId"]: {
                                "requireConfirmedCandidates": True,
                            }
                        },
                    },
                },
            )
            self.assertEqual(approve_response.status_code, 200)

            completed = self._wait_for_job_terminal_state(job_id)
            self.assertEqual(completed["status"], "blocked")
            failed_logs = [item for item in completed["logs"] if item["status"] == "failed"]
            self.assertTrue(failed_logs)
            self.assertIn("低置信度填位", failed_logs[0]["message"])

    def test_job_api_manual_mapping_can_skip_and_force_overwrite(self) -> None:
        upload_response = self.client.post(
            "/api/workflow/uploads",
            files=[
                (
                    "files",
                    (
                        "原始数据.xlsx",
                        self._build_xlsx_bytes(),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
                (
                    "files",
                    (
                        "冲突模板.xlsx",
                        self._build_label_value_template_xlsx_bytes(existing_sum=600),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ),
                ),
            ],
        )

        self.assertEqual(upload_response.status_code, 200)
        uploaded_documents = upload_response.json()["documents"]

        with patch.object(workflow_agent_service._agent_runtime, "analyze_workflow", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "summarize_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "revise_documents", return_value=None), \
             patch.object(workflow_agent_service._agent_runtime, "compose_final_output", return_value=None):
            create_response = self.client.post(
                "/api/workflow/jobs",
                json={
                    "task": "读取原始数据.xlsx 的金额，计算合计和平均值，并填到冲突模板.xlsx。",
                    "documents": uploaded_documents,
                    "manualConfirm": True,
                    "approved": False,
                    "agentMode": "off",
                    "localFallbackEnabled": True,
                    "preferredSkills": [],
                    "context": {},
                    "planOverrides": {},
                },
            )

            self.assertEqual(create_response.status_code, 200)
            job_id = create_response.json()["jobId"]
            awaiting = self._wait_for_job_terminal_state(job_id)
            table_filler_step = next(step for step in awaiting["plan"] if step.get("skill") == "table_filler")

            approve_response = self.client.post(
                f"/api/workflow/jobs/{job_id}/approve",
                json={
                    "approved": True,
                    "planOverrides": {
                        "disabledStepIds": [],
                        "stepInputOverrides": {
                            table_filler_step["stepId"]: {
                                "requireConfirmedCandidates": True,
                                "manualMappings": [
                                    {"mappingId": "map_1", "sheet": "汇总页", "cell": "B2", "writePolicy": "force_overwrite"},
                                    {"mappingId": "map_2", "enabled": False, "sheet": "汇总页", "cell": "B3"},
                                ],
                            }
                        },
                    },
                },
            )
            self.assertEqual(approve_response.status_code, 200)

            completed = self._wait_for_job_terminal_state(job_id)
            self.assertEqual(completed["status"], "completed")
            artifact_response = self.client.get(completed["finalOutput"]["downloads"][0]["downloadUrl"])
            self.assertEqual(artifact_response.status_code, 200)
            workbook = load_workbook(filename=BytesIO(artifact_response.content))
            summary_sheet = workbook["汇总页"]
            self.assertEqual(summary_sheet["B2"].value, 651.6)
            self.assertIsNone(summary_sheet["B3"].value)

            table_result = completed["finalOutput"]["artifacts"]["技能结果"]["table_filler"]
            audit_by_status = {item["status"] for item in table_result["fillAudit"]}
            audit_by_cell = {item["cell"]: item for item in table_result["fillAudit"] if item.get("cell")}
            self.assertIn("manual_mapping_skipped", audit_by_status)
            self.assertEqual(audit_by_cell["B2"]["decisionSource"], "manual")
            self.assertEqual(audit_by_cell["B2"]["writePolicy"], "force_overwrite")

    def test_sync_execute_downloads_include_artifact_metadata(self) -> None:
        response = self.client.post(
            "/api/workflow/execute",
            json={
                "task": "读取这份 Excel，计算金额合计和平均值，再按结果填表输出",
                "documents": [
                    {
                        "name": "费用明细.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 费用明细\n项目 | 金额\n打车 | 128\n酒店 | 523.6",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "费用明细",
                                    "headers": ["项目", "金额"],
                                    "rows": [
                                        {"项目": "打车", "金额": 128},
                                        {"项目": "酒店", "金额": 523.6},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["金额"],
                                }
                            ],
                        },
                    }
                ],
                "manualConfirm": False,
                "approved": True,
                "agentMode": "off",
                "localFallbackEnabled": True,
                "preferredSkills": [],
                "context": {},
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        download = payload["finalOutput"]["downloads"][0]
        self.assertTrue(download["artifactId"])
        self.assertTrue(download["downloadUrl"])
        self.assertTrue(download["contentBase64"])

    def test_job_api_executes_workflow_and_downloads_artifact(self) -> None:
        response = self.client.post(
            "/api/workflow/jobs",
            json={
                "task": "读取这份 Excel，计算金额合计和平均值，再按结果填表输出",
                "documents": [
                    {
                        "name": "费用明细.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] 费用明细\n项目 | 金额\n打车 | 128\n酒店 | 523.6",
                        "structuredData": {
                            "sheetCount": 1,
                            "sheets": [
                                {
                                    "title": "费用明细",
                                    "headers": ["项目", "金额"],
                                    "rows": [
                                        {"项目": "打车", "金额": 128},
                                        {"项目": "酒店", "金额": 523.6},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["金额"],
                                }
                            ],
                        },
                    }
                ],
                "manualConfirm": False,
                "approved": True,
                "agentMode": "off",
                "localFallbackEnabled": True,
                "preferredSkills": [],
                "context": {},
                "planOverrides": {},
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = self._wait_for_job_terminal_state(response.json()["jobId"])
        self.assertEqual(payload["status"], "completed")
        download = payload["finalOutput"]["downloads"][0]
        artifact_response = self.client.get(download["downloadUrl"])
        self.assertEqual(artifact_response.status_code, 200)
        workbook = load_workbook(filename=BytesIO(artifact_response.content))
        self.assertIn("结果表", workbook.sheetnames)

    def test_job_api_awaits_confirmation_and_approve_resumes(self) -> None:
        response = self.client.post(
            "/api/workflow/jobs",
            json={
                "task": "请帮我修订这份说明，整理得更清晰",
                "documents": [
                    {
                        "name": "draft.txt",
                        "type": "text",
                        "contentText": "原始草稿内容。",
                    }
                ],
                "manualConfirm": False,
                "approved": False,
                "agentMode": "off",
                "localFallbackEnabled": True,
                "preferredSkills": [],
                "context": {},
                "planOverrides": {},
            },
        )

        self.assertEqual(response.status_code, 200)
        job_id = response.json()["jobId"]
        awaiting = self._wait_for_job_terminal_state(job_id)
        self.assertEqual(awaiting["status"], "awaiting_confirmation")

        approve_response = self.client.post(f"/api/workflow/jobs/{job_id}/approve", json={"approved": True, "planOverrides": {}})
        self.assertEqual(approve_response.status_code, 200)
        completed = self._wait_for_job_terminal_state(job_id)
        self.assertEqual(completed["status"], "completed")
        self.assertTrue(completed["finalOutput"]["revisedDocument"])

    def test_execute_esg_outputs_topic_matrix(self) -> None:
        response = self.client.post(
            "/api/workflow/execute",
            json={
                "task": "请检查这套 ESG 披露材料的覆盖度，提取关键 KPI，并给出报告章节大纲。",
                "documents": [
                    {
                        "name": "2025_esg_report_draft.docx",
                        "type": "word",
                        "contentText": "环境部分包括可再生电力占比 48%。社会部分包括培训覆盖率 91%。治理部分包括反舞弊培训覆盖率 100%。",
                    }
                ],
                "manualConfirm": False,
                "approved": True,
                "agentMode": "off",
                "localFallbackEnabled": True,
                "preferredSkills": [],
                "context": {},
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        skill_results = payload["finalOutput"]["artifacts"]["技能结果"]
        self.assertIn("esg_standard_selector", payload["executedSkills"])
        self.assertIn("esg_disclosure_matrix_builder", payload["executedSkills"])
        self.assertIn("esg_data_request_builder", payload["executedSkills"])
        self.assertIn("esg_evidence_linker", payload["executedSkills"])
        matrix = payload["finalOutput"]["artifacts"]["技能结果"]["esg_material_checker"]["topicMatrix"]
        self.assertTrue(matrix)
        self.assertIn(matrix[0]["status"], {"covered", "weak", "missing"})
        self.assertTrue(skill_results["esg_disclosure_matrix_builder"]["disclosureMatrix"])
        self.assertTrue(skill_results["esg_data_request_builder"]["dataRequests"])
        self.assertTrue(skill_results["esg_evidence_linker"]["evidenceLinks"])

    def test_execute_spreadsheet_workflow_supports_cross_sheet_ratio_and_trend(self) -> None:
        response = self.client.post(
            "/api/workflow/execute",
            json={
                "task": "请跨工作表汇总 Scope 1 和 Scope 2，计算占比、同比和环比。",
                "documents": [
                    {
                        "name": "emissions.xlsx",
                        "type": "excel",
                        "contentText": "[工作表] Jan\n月份 | Scope 1 | Scope 2\n1月 | 100 | 200\n2月 | 110 | 220\n\n[工作表] Feb\n月份 | Scope 1 | Scope 2\n3月 | 120 | 240\n4月 | 150 | 260",
                        "structuredData": {
                            "sheetCount": 2,
                            "sheets": [
                                {
                                    "title": "Jan",
                                    "headers": ["月份", "Scope 1", "Scope 2"],
                                    "rows": [
                                        {"月份": "1月", "Scope 1": 100, "Scope 2": 200},
                                        {"月份": "2月", "Scope 1": 110, "Scope 2": 220},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["Scope 1", "Scope 2"],
                                },
                                {
                                    "title": "Feb",
                                    "headers": ["月份", "Scope 1", "Scope 2"],
                                    "rows": [
                                        {"月份": "3月", "Scope 1": 120, "Scope 2": 240},
                                        {"月份": "4月", "Scope 1": 150, "Scope 2": 260},
                                    ],
                                    "rowCount": 2,
                                    "numericColumns": ["Scope 1", "Scope 2"],
                                },
                            ],
                        },
                    }
                ],
                "manualConfirm": False,
                "approved": True,
                "agentMode": "off",
                "localFallbackEnabled": True,
                "preferredSkills": [],
                "context": {},
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        results = payload["finalOutput"]["artifacts"]["技能结果"]["spreadsheet_calculator"]["results"]
        cross_sheet = [item for item in results if item["sheet"] == "跨工作表"]
        self.assertTrue(cross_sheet)
        self.assertIn("ratio", cross_sheet[0]["metrics"])

    def test_prepare_documents_merges_content_text_and_ocr_text(self) -> None:
        prepared = WorkflowInputService().prepare_documents(
            [
                WorkflowDocument(
                    name="scanned_esg_report.pdf",
                    type="pdf",
                    contentText="治理章节：董事会 ESG 委员会召开 4 次。",
                    ocrText="治理章节：董事会 ESG 委员会召开 4 次。\n图表说明：反舞弊培训覆盖率 100%。",
                )
            ]
        )

        self.assertEqual(len(prepared), 1)
        document = prepared[0]
        self.assertTrue(document.usedOcr)
        self.assertIn("反舞弊培训覆盖率 100%", document.text)
        self.assertIn("治理章节：董事会 ESG 委员会召开 4 次。", document.text)


if __name__ == "__main__":
    unittest.main()
